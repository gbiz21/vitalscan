"""Inject screenshot-capture instructions into slide 6 of the team's deck.

This script edits the existing .pptx in place so the design Jason built is
preserved — only the speaker notes on slide 6 change.
"""
from pathlib import Path

from pptx import Presentation

DECK = Path(__file__).resolve().parent / "VitalScan_Group1_Backend_System_Architecture.pptx"

SLIDE_6_NOTES = """LIVE SITE (90s · Germaine) — and how to capture the 4 screenshots before presenting:

═══════════════════════════════════════════════════════════════
WHAT TO SAY (90s)
═══════════════════════════════════════════════════════════════

I want to step through what the site actually does in four images.

Image 1 — the webcam capture flow. The user lands on vitalscan.bkre8tive.com,
clicks "Scan live", grants camera permission, and the browser records for
30 seconds with a visible countdown and recording indicator.

Image 2 — the dashboard once the scan completes. A body figure on the left
with marker positions for each biomarker, and four metric cards on the right.

Image 3 — zoom on the biomarker cards. Heart rate in BPM, HRV in milliseconds,
stress index between 0 and 1, and blood pressure. The first three come from
the rPPG pipeline. Blood pressure is mocked and labeled as a known limit.

Image 4 — for the engineers in the room, the /scan response JSON visible in
the browser's network tab. This is the shared API contract any other team
can consume.

The site is live right now — open vitalscan.bkre8tive.com on your phone
and try it.

═══════════════════════════════════════════════════════════════
HOW TO CAPTURE THE 4 SCREENSHOTS (do this BEFORE the talk)
═══════════════════════════════════════════════════════════════

Browser: Chrome on macOS. Files save to ~/Desktop by default.

──────────────────────────────────────────────────────────────
SCREENSHOT 1 — Webcam capture (30-second countdown)
──────────────────────────────────────────────────────────────
1. Open https://vitalscan.bkre8tive.com in Chrome
2. Click the "Scan live" button
3. Allow camera permission when prompted
4. WAIT for recording to start — you'll see corner brackets,
   a sweeping scan line, and a red "REC" indicator
5. While recording is active, press Cmd+Shift+5
6. Click "Capture Selected Window" → click the Chrome window
7. Done — capture goes to Desktop as a PNG

──────────────────────────────────────────────────────────────
SCREENSHOT 2 — Dashboard (body figure + biomarker cards)
──────────────────────────────────────────────────────────────
1. After a scan completes, you'll see the body figure + 4 cards
   (a persisted prior scan from a refresh is fine for this shot)
2. Zoom browser to 100% (Cmd+0) so the layout is clean
3. Press Cmd+Shift+4 (cursor turns to crosshair)
4. Drag a tight rectangle around the whole panel — header,
   body figure on the left, all 4 metric cards on the right
5. Release — capture saved to Desktop

──────────────────────────────────────────────────────────────
SCREENSHOT 3 — Biomarker cards close-up (HR · HRV · stress · BP)
──────────────────────────────────────────────────────────────
1. Same dashboard view as Screenshot 2
2. Press Cmd-Plus 2–3 times to zoom in so the numbers are large
3. Press Cmd+Shift+4
4. Drag around JUST the 4 biomarker cards (no body figure)
5. The four big numbers (80 bpm, 42 ms, 0.52, 126/84) should
   be unmistakable in this shot — that's the "money shot"
6. Cmd+0 afterward to reset zoom

──────────────────────────────────────────────────────────────
SCREENSHOT 4 — JSON response (developer view)
──────────────────────────────────────────────────────────────
1. Open Chrome DevTools: Cmd+Option+I
2. Click the "Network" tab in DevTools
3. In the network filter bar, type:  scan
4. Click "Scan live" or "Rescan" on the site to trigger a request
5. When the /scan request appears in the list, click it
6. In the right pane, click the "Response" tab
7. You should see the JSON:
       { "biomarkers": {
           "heart_rate": 80,
           "hrv_sdnn": 42,
           "stress_index": 0.52,
           "blood_pressure": { "systolic": 126, "diastolic": 84 }
       } }
8. Press Cmd+Shift+4 and drag around the DevTools panel
   showing BOTH the highlighted /scan row AND the JSON response

═══════════════════════════════════════════════════════════════
HOW TO DROP THE SCREENSHOTS INTO THIS SLIDE
═══════════════════════════════════════════════════════════════
Option A — drag & drop:
  In PowerPoint / Keynote, open this deck. Drag each PNG from
  Finder directly onto its matching placeholder rectangle on
  slide 6. Resize so it fills the box.

Option B — right-click replace:
  Right-click a placeholder → "Change Picture" → "From File"
  → select the screenshot. PowerPoint preserves the placeholder
  size automatically.

After all 4 images are in, delete the gray "Drop screenshot here"
labels if PowerPoint didn't replace them automatically.
"""


def main() -> None:
    prs = Presentation(DECK)
    if len(prs.slides) < 6:
        raise SystemExit(f"deck has only {len(prs.slides)} slides — expected ≥6")

    slide6 = prs.slides[5]
    slide6.notes_slide.notes_text_frame.text = SLIDE_6_NOTES

    prs.save(DECK)
    print(f"updated speaker notes on slide 6 of {DECK.name}")
    print(f"notes length: {len(SLIDE_6_NOTES)} chars")


if __name__ == "__main__":
    main()
