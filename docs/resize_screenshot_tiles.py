"""Resize slide 6 tiles AND fit any inserted pictures into them.

Original tile layout:
    Outer card    T=1.65  H=2.20
    Inner image   T=1.80  H=1.25   ← squashes screenshots (2.1:1 aspect)
    "Drop here"   T=2.32
    Title         T=3.27
    Caption       T=3.57
    (gap)
    Bottom proof  T=5.15

Target tile layout (grows into the unused space above the proof callout):
    Outer card    T=1.65  H=3.30
    Inner image   T=1.80  H=1.95   ← aspect ~4:3, screenshot-natural
    "Drop here"   T=2.68
    Title         T=3.95
    Caption       T=4.25
    Bottom proof  T=5.15   (unchanged)

For each inserted picture, fits it into its tile's inner zone preserving the
picture's aspect ratio (letterbox/pillarbox rather than distort).

Run after closing PowerPoint:
    backend/venv312/bin/python docs/resize_screenshot_tiles.py
"""
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Emu, Inches

DECK = Path(__file__).resolve().parent / "VitalScan_Group1_Backend_System_Architecture.pptx"

TOL = Emu(50_000)  # ~0.05" tolerance when matching positions


def near(a, b):
    return abs(a - b) < TOL


def fit_picture(pic, box_L, box_T, box_W, box_H):
    """Resize pic to fit inside the box, preserving its aspect ratio."""
    pic_aspect = pic.width / pic.height
    box_aspect = box_W / box_H
    if pic_aspect > box_aspect:
        # picture is wider than box → fit by width, letterbox top/bottom
        new_W = box_W
        new_H = int(box_W / pic_aspect)
    else:
        # picture is taller (or equal) → fit by height, pillarbox left/right
        new_H = box_H
        new_W = int(box_H * pic_aspect)
    pic.width = new_W
    pic.height = new_H
    pic.left = box_L + (box_W - new_W) // 2
    pic.top = box_T + (box_H - new_H) // 2


def main():
    prs = Presentation(DECK)
    slide = prs.slides[5]

    new_inner_zones: list[tuple[int, int, int, int]] = []

    # PASS 1: regrow the tile rectangles and move the labels.
    for sh in slide.shapes:
        T, H = sh.top, sh.height

        # OUTER tile (handles both already-grown and original states)
        if near(T, Inches(1.65)) and (near(H, Inches(2.20)) or near(H, Inches(3.30))):
            sh.height = Inches(3.30)

        # INNER screenshot zone
        elif near(T, Inches(1.80)) and (near(H, Inches(1.25)) or near(H, Inches(1.95))):
            sh.height = Inches(1.95)
            new_inner_zones.append((sh.left, sh.top, sh.width, sh.height))

        # "Drop screenshot here" labels — re-center vertically
        elif near(T, Inches(2.32)) or near(T, Inches(2.68)):
            sh.top = Inches(2.68)

        # Title labels
        elif near(T, Inches(3.27)) or near(T, Inches(3.95)):
            sh.top = Inches(3.95)

        # Caption labels
        elif near(T, Inches(3.57)) or near(T, Inches(4.25)):
            sh.top = Inches(4.25)

    new_inner_zones.sort()  # left-to-right

    # PASS 2: snap any inserted pictures into their tile zones.
    pics = [sh for sh in slide.shapes if sh.shape_type == MSO_SHAPE_TYPE.PICTURE]
    pics.sort(key=lambda p: p.left)

    pics_fit = 0
    for pic in pics:
        cx = pic.left + pic.width // 2
        target = next(
            (zone for zone in new_inner_zones
             if zone[0] <= cx <= zone[0] + zone[2]),
            None,
        )
        if target is None:
            print(f"warning: picture at L={pic.left / 914400:.2f}\" "
                  f"doesn't fall inside any inner zone — skipped")
            continue
        fit_picture(pic, *target)
        pics_fit += 1

    prs.save(DECK)
    print(f"resized slide 6 in {DECK.name}")
    print(f"  tile rectangles regrown: 4 outer + 4 inner")
    print(f"  labels repositioned: 12 (4 placeholders × 3 lines)")
    print(f"  pictures fit-to-zone: {pics_fit}/{len(pics)}")


if __name__ == "__main__":
    main()
