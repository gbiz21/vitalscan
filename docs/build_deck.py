"""Build the AIT 500 Group 1 (VitalScan) Week-8 presentation deck.

Per the 2026-05-13 meeting:
  • 10–15 minutes total · ~1 slide per group member
  • Show progress, not just final results
  • Include site screenshots + dataset retrieval story
  • Group 1 members: Germaine, Jason, Daray, Abner

The architecture slide is built from native PPTX shapes so each block can fade
in on click — animations are injected as raw <p:timing> XML.
"""
from copy import deepcopy
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import nsmap, qn
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parent.parent
OUT = ROOT / "docs" / "VitalScan_Group1_Presentation.pptx"

# --- palette (matches the figures) ---
NAVY = RGBColor(0x0F, 0x1F, 0x44)
SLATE = RGBColor(0x33, 0x44, 0x5C)
INK = RGBColor(0x1A, 0x1A, 0x1A)
MUTED = RGBColor(0x55, 0x66, 0x7A)
ACCENT = RGBColor(0x1F, 0x6F, 0xEB)
GOOD = RGBColor(0x1E, 0x88, 0x4F)
WARN = RGBColor(0xC0, 0x73, 0x1D)
SOFT_BG = RGBColor(0xF6, 0xF8, 0xFB)
CARD_BG = RGBColor(0xEE, 0xF3, 0xFA)
GREEN_BG = RGBColor(0xE9, 0xF5, 0xEC)
ORANGE_BG = RGBColor(0xFD, 0xF1, 0xE3)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


# ---------- helpers ----------
def add_slide():
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE
    return s


def add_text(slide, left, top, width, height, text, *, size=18, bold=False,
             color=INK, align=PP_ALIGN.LEFT, font="Calibri"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
    return tb


def add_bullets(slide, left, top, width, height, bullets, *, size=18, color=INK):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(6)
        r = p.add_run()
        r.text = f"•  {b}"
        r.font.name = "Calibri"
        r.font.size = Pt(size)
        r.font.color.rgb = color
    return tb


def header(slide, title, subtitle=None, presenter=None):
    stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, Inches(0.18))
    stripe.line.fill.background()
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = ACCENT
    add_text(slide, Inches(0.5), Inches(0.35), Inches(12.3), Inches(0.7),
             title, size=28, bold=True, color=NAVY)
    if subtitle:
        add_text(slide, Inches(0.5), Inches(0.95), Inches(12.3), Inches(0.4),
                 subtitle, size=14, color=MUTED)
    add_text(slide, Inches(0.5), Inches(7.05), Inches(8.0), Inches(0.3),
             "VitalScan · AIT 500 · Group 1 · Week 8",
             size=10, color=MUTED)
    if presenter:
        add_text(slide, Inches(8.5), Inches(7.05), Inches(4.3), Inches(0.3),
                 f"presented by {presenter}", size=10, color=ACCENT,
                 align=PP_ALIGN.RIGHT)


def set_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def filled_card(slide, left, top, width, height, *,
                fill=SOFT_BG, border=ACCENT, radius=0.06):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  left, top, width, height)
    card.adjustments[0] = radius
    card.fill.solid()
    card.fill.fore_color.rgb = fill
    card.line.color.rgb = border
    card.line.width = Pt(1.25)
    card.shadow.inherit = False
    return card


# ---------- animation helper ----------
def add_sequential_fade(slide, groups, *, dur_ms=600):
    """Inject <p:timing> so each group of shapes fades in on click.

    groups: list of lists of shape ids.  One click reveals one inner list.
    """
    p = "{%s}" % P_NS
    timing = etree.SubElement(slide._element, f"{p}timing")
    tnLst = etree.SubElement(timing, f"{p}tnLst")
    root_par = etree.SubElement(tnLst, f"{p}par")
    root_ctn = etree.SubElement(
        root_par, f"{p}cTn",
        id="1", dur="indefinite", restart="never", nodeType="tmRoot",
    )
    root_children = etree.SubElement(root_ctn, f"{p}childTnLst")
    seq = etree.SubElement(root_children, f"{p}seq",
                           concurrent="1", nextAc="seek")
    main_ctn = etree.SubElement(seq, f"{p}cTn",
                                id="2", dur="indefinite", nodeType="mainSeq")
    main_children = etree.SubElement(main_ctn, f"{p}childTnLst")

    # cond lists at end of seq
    prev_cond = etree.SubElement(seq, f"{p}prevCondLst")
    cond1 = etree.SubElement(prev_cond, f"{p}cond", evt="onPrev", delay="0")
    tgt1 = etree.SubElement(cond1, f"{p}tgtEl")
    etree.SubElement(tgt1, f"{p}sldTgt")
    next_cond = etree.SubElement(seq, f"{p}nextCondLst")
    cond2 = etree.SubElement(next_cond, f"{p}cond", evt="onNext", delay="0")
    tgt2 = etree.SubElement(cond2, f"{p}tgtEl")
    etree.SubElement(tgt2, f"{p}sldTgt")

    next_id = [3]

    def nid():
        i = next_id[0]
        next_id[0] += 1
        return str(i)

    for group_idx, shape_ids in enumerate(groups):
        click_par = etree.SubElement(main_children, f"{p}par")
        click_ctn = etree.SubElement(click_par, f"{p}cTn", id=nid(), fill="hold")
        click_stcond = etree.SubElement(click_ctn, f"{p}stCondLst")
        etree.SubElement(click_stcond, f"{p}cond", delay="indefinite")
        click_children = etree.SubElement(click_ctn, f"{p}childTnLst")

        grp_par = etree.SubElement(click_children, f"{p}par")
        grp_ctn = etree.SubElement(grp_par, f"{p}cTn", id=nid(), fill="hold")
        grp_stcond = etree.SubElement(grp_ctn, f"{p}stCondLst")
        etree.SubElement(grp_stcond, f"{p}cond", delay="0")
        grp_children = etree.SubElement(grp_ctn, f"{p}childTnLst")

        for shape_idx, sid in enumerate(shape_ids):
            # one effect per shape; only the first is a click effect, others "withEffect"
            node_type = "clickEffect" if shape_idx == 0 else "withEffect"
            effect_par = etree.SubElement(grp_children, f"{p}par")
            effect_ctn = etree.SubElement(
                effect_par, f"{p}cTn",
                id=nid(), presetID="10", presetClass="entr",
                presetSubtype="0", fill="hold", grpId="0",
                nodeType=node_type,
            )
            effect_stcond = etree.SubElement(effect_ctn, f"{p}stCondLst")
            etree.SubElement(effect_stcond, f"{p}cond", delay="0")
            effect_children = etree.SubElement(effect_ctn, f"{p}childTnLst")

            # visibility set
            set_el = etree.SubElement(effect_children, f"{p}set")
            cbhvr = etree.SubElement(set_el, f"{p}cBhvr")
            etree.SubElement(cbhvr, f"{p}cTn", id=nid(), dur="1", fill="hold")
            tgt = etree.SubElement(cbhvr, f"{p}tgtEl")
            etree.SubElement(tgt, f"{p}spTgt", spid=str(sid))
            attrlst = etree.SubElement(cbhvr, f"{p}attrNameLst")
            an = etree.SubElement(attrlst, f"{p}attrName")
            an.text = "style.visibility"
            to_el = etree.SubElement(set_el, f"{p}to")
            etree.SubElement(to_el, f"{p}strVal", val="visible")

            # opacity tween
            anim = etree.SubElement(effect_children, f"{p}anim",
                                    calcmode="lin", valueType="num")
            anim_cbhvr = etree.SubElement(anim, f"{p}cBhvr", additive="base")
            etree.SubElement(anim_cbhvr, f"{p}cTn",
                             id=nid(), dur=str(dur_ms), fill="hold")
            anim_tgt = etree.SubElement(anim_cbhvr, f"{p}tgtEl")
            etree.SubElement(anim_tgt, f"{p}spTgt", spid=str(sid))
            anim_attrlst = etree.SubElement(anim_cbhvr, f"{p}attrNameLst")
            anim_an = etree.SubElement(anim_attrlst, f"{p}attrName")
            anim_an.text = "style.opacity"
            tavLst = etree.SubElement(anim, f"{p}tavLst")
            tav0 = etree.SubElement(tavLst, f"{p}tav", tm="0")
            val0 = etree.SubElement(tav0, f"{p}val")
            etree.SubElement(val0, f"{p}fltVal", val="0")
            tav1 = etree.SubElement(tavLst, f"{p}tav", tm="100000")
            val1 = etree.SubElement(tav1, f"{p}val")
            etree.SubElement(val1, f"{p}fltVal", val="1")


# ============================================================
# SLIDE 1 — TITLE / MEMBERS
# ============================================================
s = add_slide()
panel = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(5.6), SH)
panel.line.fill.background()
panel.fill.solid()
panel.fill.fore_color.rgb = NAVY

add_text(s, Inches(0.6), Inches(0.6), Inches(4.8), Inches(0.4),
         "AIT 500 · WESTCLIFF UNIVERSITY · WEEK 8",
         size=11, bold=True, color=ACCENT)
add_text(s, Inches(0.6), Inches(1.3), Inches(5.0), Inches(1.6),
         "VitalScan", size=66, bold=True, color=WHITE)
add_text(s, Inches(0.6), Inches(2.8), Inches(5.0), Inches(0.8),
         "Vital signs from a 30-second\nwebcam video",
         size=20, color=WHITE)
add_text(s, Inches(0.6), Inches(4.3), Inches(4.8), Inches(0.4),
         "GROUP 1 · rPPG SIGNAL EXTRACTION",
         size=12, bold=True, color=ACCENT)
add_text(s, Inches(0.6), Inches(6.6), Inches(4.8), Inches(0.4),
         "vitalscan.bkre8tive.com", size=14, bold=True, color=ACCENT)

# right side — members + status
add_text(s, Inches(6.2), Inches(0.9), Inches(6.5), Inches(0.4),
         "GROUP 1 — TASK 1: rPPG SIGNAL EXTRACTION",
         size=12, bold=True, color=MUTED)
add_text(s, Inches(6.2), Inches(1.35), Inches(6.5), Inches(0.7),
         "Project status & progress",
         size=30, bold=True, color=NAVY)

# member cards
def member_card(left, top, name, role):
    card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                              left, top, Inches(6.3), Inches(0.85))
    card.adjustments[0] = 0.18
    card.fill.solid(); card.fill.fore_color.rgb = CARD_BG
    card.line.color.rgb = ACCENT; card.line.width = Pt(1)
    add_text(s, left + Inches(0.25), top + Inches(0.13),
             Inches(2.4), Inches(0.4), name, size=15, bold=True, color=NAVY)
    add_text(s, left + Inches(2.7), top + Inches(0.17),
             Inches(3.6), Inches(0.4), role, size=12, color=SLATE)

member_card(Inches(6.2), Inches(2.55), "Germaine Beazer",   "Full-stack site · deployment")
member_card(Inches(6.2), Inches(3.50), "Jason",             "Dataset retrieval · slide draft")
member_card(Inches(6.2), Inches(4.45), "Daray",             "rPPG signal · face detection / ROI")
member_card(Inches(6.2), Inches(5.40), "Abner",             "Task share TBD")

add_text(s, Inches(6.2), Inches(6.55), Inches(6.5), Inches(0.4),
         "github.com/gbiz21/vitalscan", size=12, color=MUTED)

set_notes(s, """OPENER (60s · Germaine, then around the team):

Good afternoon. We are Group 1 for the VitalScan project. Our assignment from the project brief was Task 1: rPPG signal extraction — face detection and region-of-interest processing for remote photoplethysmography.

We're going to spend about ten to fifteen minutes today walking you through what we built so far, what data sets we pulled in, and where the project stands. The site is live at vitalscan.bkre8tive.com — you can open it on your phone while we're presenting.

Quick round of introductions before we start so everyone knows who is talking — Germaine has been driving the build and the deployment. Jason coordinated the dataset retrieval. Daray's been on the rPPG signal side — face detection and ROIs. Abner — [if present: introduce; if not: note that Abner couldn't make today's session, but their part of the writeup is included].

Each of us will take one or two slides.""")

# ============================================================
# SLIDE 2 — THE ASSIGNMENT
# ============================================================
s = add_slide()
header(s, "The assignment", "What the project brief asked Group 1 to deliver",
       presenter="Germaine")

# left — brief excerpt
filled_card(s, Inches(0.5), Inches(1.6), Inches(6.0), Inches(4.8))
add_text(s, Inches(0.75), Inches(1.8), Inches(5.5), Inches(0.4),
         "GROUP 1 · rPPG SIGNAL EXTRACTION", size=12, bold=True, color=ACCENT)
add_bullets(s, Inches(0.75), Inches(2.3), Inches(5.5), Inches(4.0), [
    "Face detection on each video frame",
    "Region-of-interest (ROI) bounding boxes on the face",
    "Extract a remote-PPG signal from those ROIs",
    "Output: heart rate, HRV, stress index, blood pressure",
    "Match the project's shared JSON contract",
], size=14)

# right — what we delivered
filled_card(s, Inches(6.8), Inches(1.6), Inches(6.0), Inches(4.8),
            fill=GREEN_BG, border=GOOD)
add_text(s, Inches(7.05), Inches(1.8), Inches(5.5), Inches(0.4),
         "WHAT WE BUILT", size=12, bold=True, color=GOOD)
add_bullets(s, Inches(7.05), Inches(2.3), Inches(5.5), Inches(4.0), [
    "Full Python rPPG pipeline (face + POS + HRV)",
    "FastAPI /scan endpoint exposing the pipeline",
    "React dashboard that calls /scan and renders biomarkers",
    "Webcam capture in-browser (no plugins)",
    "Deployed live — anyone in this room can hit it now",
], size=14)

add_text(s, Inches(0.5), Inches(6.65), Inches(12.3), Inches(0.4),
         "We took Task 1 (the brief's minimum) and built it out to a full working product to make our progress easy to demo.",
         size=12, color=MUTED, align=PP_ALIGN.CENTER)

set_notes(s, """THE ASSIGNMENT (90s · Germaine):

Project brief, in plain words: there were four groups, and each group owned one slice of the pipeline. Our slice — Group 1 — was rPPG signal extraction. Face detection. ROI bounding boxes. Pull the heart-rate signal out of the video. Format the output as the JSON the other groups expect.

That's the left column — the minimum we were asked for.

The right column is what we actually shipped. We didn't stop at "produce a JSON file." We built the Python pipeline, wrapped it in a REST API, put a React dashboard on top that calls the API, added browser webcam capture so you don't need to upload a file, and deployed the whole thing to a live URL on my home server. The site is at vitalscan.bkre8tive.com — TLS, no port forwarding, runs behind a Cloudflare tunnel.

We did this so the demo is easy. Instead of showing the professor a console output, we can hand him a phone.""")

# ============================================================
# SLIDE 3 — DATASETS (Jason)
# ============================================================
s = add_slide()
header(s, "Datasets retrieved",
       "Pulled per the project brief — Google Drive + Microsoft sources",
       presenter="Jason")

# UBFC card
ubfc = filled_card(s, Inches(0.5), Inches(1.65), Inches(6.0), Inches(4.6),
                   fill=CARD_BG, border=ACCENT)
add_text(s, Inches(0.75), Inches(1.85), Inches(5.5), Inches(0.5),
         "UBFC-rPPG", size=22, bold=True, color=NAVY)
add_text(s, Inches(0.75), Inches(2.35), Inches(5.5), Inches(0.4),
         "Real human subjects · Univ. de Bourgogne, France",
         size=12, color=MUTED)
add_bullets(s, Inches(0.75), Inches(2.85), Inches(5.5), Inches(3.2), [
    "42 subjects · webcam at 30 fps",
    "Synchronized finger pulse-oximeter ground truth",
    "Source: Google Drive (linked in project brief)",
    "Headline dataset — this is what we evaluate against",
], size=13)

# SCAMPS card
scamps = filled_card(s, Inches(6.8), Inches(1.65), Inches(6.0), Inches(4.6),
                     fill=ORANGE_BG, border=WARN)
add_text(s, Inches(7.05), Inches(1.85), Inches(5.5), Inches(0.5),
         "SCAMPS", size=22, bold=True, color=NAVY)
add_text(s, Inches(7.05), Inches(2.35), Inches(5.5), Inches(0.4),
         "Synthetic avatars · Microsoft Research",
         size=12, color=MUTED)
add_bullets(s, Inches(7.05), Inches(2.85), Inches(5.5), Inches(3.2), [
    "10 example videos (public subset on GitHub)",
    "Simulated PPG ground truth from the avatar",
    "Source: Microsoft (danmcduff/scampsdataset)",
    "Used as a synthetic-data baseline",
], size=13)

add_text(s, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.5),
         "Both datasets are embedded in our pipeline · same code reads both · evaluation report in docs/evaluation_results.md",
         size=12, color=MUTED, align=PP_ALIGN.CENTER)

set_notes(s, """DATASETS (90s · Jason):

The brief told us to pull two datasets for evaluation, one real and one synthetic. We retrieved both.

UBFC-rPPG, on the left, is the dataset from the University of Burgundy. 42 real human subjects sitting in front of a webcam, with a synchronized finger pulse oximeter recording ground-truth heart rate at the same time. The download link was on the project's Google Drive. This is the dataset that actually matters for our use case because the end product is for real people.

SCAMPS, on the right, is Microsoft's synthetic-avatar dataset. Computer-generated faces with a simulated pulse signal. The full SCAMPS dataset lives in Microsoft Azure storage that we don't have access to, but Microsoft also published a 10-video example subset on GitHub, which is what we used. It serves as a synthetic-data baseline — useful for stress-testing the algorithm.

Both datasets are read by the same code in our evaluation harness, so the comparison is apples to apples. The full numbers are in our evaluation report on the repo. The short version is on the next-next slide.""")

# ============================================================
# SLIDE 4 — rPPG APPROACH (Daray)
# ============================================================
s = add_slide()
header(s, "rPPG signal extraction — how it works",
       "Face → ROI → green-channel averaging → POS → heart rate",
       presenter="Daray")

# 4-step pipeline cards
def step_card(left, num, title, body):
    card = filled_card(s, left, Inches(1.65), Inches(2.9), Inches(4.7))
    # number circle
    circle = s.shapes.add_shape(MSO_SHAPE.OVAL,
                               left + Inches(0.25), Inches(1.85),
                               Inches(0.7), Inches(0.7))
    circle.fill.solid(); circle.fill.fore_color.rgb = ACCENT
    circle.line.fill.background()
    add_text(s, left + Inches(0.25), Inches(1.9), Inches(0.7), Inches(0.6),
             str(num), size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, left + Inches(0.25), Inches(2.7), Inches(2.4), Inches(0.6),
             title, size=15, bold=True, color=NAVY)
    add_bullets(s, left + Inches(0.25), Inches(3.4), Inches(2.4), Inches(2.8),
                body, size=11)

step_card(Inches(0.5), 1, "Detect face",
          ["MediaPipe FaceMesh",
           "468 landmarks / frame",
           "Real-time on CPU"])
step_card(Inches(3.6), 2, "Pick ROIs",
          ["Forehead + 2 cheeks",
           "Highest capillary density",
           "Lowest motion artifact"])
step_card(Inches(6.7), 3, "POS algorithm",
          ["Project RGB onto plane",
           "orthogonal to skin tone",
           "Pulse ↑, noise ↓"])
step_card(Inches(9.8), 4, "FFT → BPM",
          ["Bandpass 0.7–4 Hz",
           "FFT, search 60–210 BPM",
           "Peak × 60 = heart rate"])

# arrows
for x in [Inches(3.35), Inches(6.45), Inches(9.55)]:
    a = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x, Inches(3.85),
                          Inches(0.3), Inches(0.4))
    a.fill.solid(); a.fill.fore_color.rgb = MUTED
    a.line.fill.background()

add_text(s, Inches(0.5), Inches(6.55), Inches(12.3), Inches(0.4),
         "Same pulse waveform is reused to compute HRV (SDNN) and a stress index from the LF/HF power ratio.",
         size=12, color=MUTED, align=PP_ALIGN.CENTER)

set_notes(s, """rPPG APPROACH (2m · Daray):

Four steps. Read left to right.

Step one — face detection. We use MediaPipe FaceMesh, which is Google's open-source face-landmark model. It puts 468 landmarks on every face in the video at real-time speed on a regular CPU.

Step two — pick the right regions. We use the forehead and both cheeks. Two reasons: those areas have the most surface capillaries, so the rPPG signal is strongest, and they move the least compared to the mouth and eyes, so they're cleaner.

Step three — POS, which stands for Plane-Orthogonal-to-Skin. The trick is that ambient lighting noise moves all three RGB channels together, but the pulse signal is mostly in the green channel. If we project the data onto the plane perpendicular to the skin-tone vector, the pulse stays and most of the noise drops out.

Step four — read the heart rate. We bandpass filter the cleaned signal between 0.7 and 4 hertz (which is 42 to 240 BPM, the physiological range), take an FFT, and pick the largest peak between 60 and 210 BPM. Multiply by 60 and you have heart rate.

We use the same cleaned signal to compute heart-rate variability — that's the standard deviation of the gaps between heartbeats — and a stress index from the low-frequency over high-frequency power ratio. Those last two are bonus biomarkers on top of the heart rate the brief asked for.""")

# ============================================================
# SLIDE 5 — ARCHITECTURE (ANIMATED) (Germaine)
# ============================================================
s = add_slide()
header(s, "System architecture — what runs where",
       "Click to step through · client → Cloudflare → homelab → eval",
       presenter="Germaine")

# track shape IDs in the order we want them revealed
groups = []   # list of lists of shape_ids

def card_with_text(left, top, w, h, title, lines, *,
                   fill=CARD_BG, border=ACCENT, title_color=NAVY,
                   title_size=14, body_size=11):
    ids = []
    card = filled_card(s, left, top, w, h, fill=fill, border=border)
    ids.append(card.shape_id)
    tb1 = add_text(s, left + Inches(0.2), top + Inches(0.15),
                   w - Inches(0.4), Inches(0.4),
                   title, size=title_size, bold=True, color=title_color)
    ids.append(tb1.shape_id)
    tb2 = add_bullets(s, left + Inches(0.2), top + Inches(0.6),
                      w - Inches(0.4), h - Inches(0.7),
                      lines, size=body_size, color=INK)
    ids.append(tb2.shape_id)
    return ids


def label(left, top, w, h, text, **kwargs):
    tb = add_text(s, left, top, w, h, text, **kwargs)
    return tb.shape_id


def arrow(left, top, w, h, *, vertical=False, color=ACCENT):
    shape = MSO_SHAPE.DOWN_ARROW if vertical else MSO_SHAPE.RIGHT_ARROW
    a = s.shapes.add_shape(shape, left, top, w, h)
    a.fill.solid(); a.fill.fore_color.rgb = color
    a.line.fill.background()
    return a.shape_id

# --- group 1: CLIENT ---
client_ids = []
section_lbl = label(Inches(0.5), Inches(1.4), Inches(12.3), Inches(0.3),
                    "STEP 1 · CLIENT (any device)", size=10,
                    bold=True, color=MUTED)
client_ids.append(section_lbl)
client_ids += card_with_text(Inches(0.6), Inches(1.75), Inches(4.0), Inches(1.45),
                             "Webcam capture",
                             ["getUserMedia + MediaRecorder",
                              "30 s WebM @ 30 fps · in-browser"])
client_ids += card_with_text(Inches(4.75), Inches(1.75), Inches(4.0), Inches(1.45),
                             "User browser",
                             ["React 18 · Vite · TypeScript",
                              "Dashboard + biomarker cards"])
client_ids += card_with_text(Inches(8.9), Inches(1.75), Inches(3.9), Inches(1.45),
                             "File upload",
                             ["<input type=\"file\">",
                              ".mp4 · .webm · .mov · ≤ 100 MB"])
groups.append(client_ids)

# --- group 2: ARROW + CLOUDFLARE ---
cf_ids = []
arrow_down1 = arrow(Inches(6.4), Inches(3.27), Inches(0.45), Inches(0.45),
                    vertical=True, color=ACCENT)
cf_ids.append(arrow_down1)
post_lbl = label(Inches(6.95), Inches(3.35), Inches(5.0), Inches(0.3),
                 "HTTPS POST /scan · multipart video", size=10,
                 color=MUTED)
cf_ids.append(post_lbl)
cf_section = label(Inches(0.5), Inches(3.85), Inches(12.3), Inches(0.3),
                   "STEP 2 · CLOUDFLARE (managed edge + tunnel)",
                   size=10, bold=True, color=MUTED)
cf_ids.append(cf_section)
cf_ids += card_with_text(Inches(0.6), Inches(4.2), Inches(6.0), Inches(1.55),
                         "Cloudflare edge",
                         ["TLS termination · DDoS · WAF",
                          "No origin certificate on homelab"])
cf_ids += card_with_text(Inches(6.85), Inches(4.2), Inches(5.95), Inches(1.55),
                         "Cloudflare Tunnel",
                         ["Named tunnel · outbound-only QUIC",
                          "No public IP · firewall stays closed"])
groups.append(cf_ids)

# --- group 3: ARROW + HOMELAB STACK ---
home_ids = []
arrow_down2 = arrow(Inches(6.4), Inches(5.82), Inches(0.45), Inches(0.45),
                    vertical=True, color=ACCENT)
home_ids.append(arrow_down2)
hl_section = label(Inches(0.5), Inches(6.35), Inches(12.3), Inches(0.3),
                   "STEP 3 · HOMELAB HOST · docker-compose",
                   size=10, bold=True, color=MUTED)
home_ids.append(hl_section)

home_ids += card_with_text(Inches(0.6), Inches(6.7), Inches(2.9), Inches(0.3),
                           "cloudflared", [],
                           fill=CARD_BG, border=ACCENT,
                           title_size=11)
home_ids += card_with_text(Inches(3.6), Inches(6.7), Inches(2.9), Inches(0.3),
                           "frontend (nginx)", [],
                           fill=GREEN_BG, border=GOOD,
                           title_size=11, title_color=GOOD)
home_ids += card_with_text(Inches(6.6), Inches(6.7), Inches(2.9), Inches(0.3),
                           "backend (FastAPI + rPPG)", [],
                           fill=GREEN_BG, border=GOOD,
                           title_size=11, title_color=GOOD)
home_ids += card_with_text(Inches(9.6), Inches(6.7), Inches(3.2), Inches(0.3),
                           "offline eval (Task 4)", [],
                           fill=ORANGE_BG, border=WARN,
                           title_size=11, title_color=WARN)
groups.append(home_ids)

# initial visibility: make all groups start hidden by setting invisible fills?
# Actually, with PowerPoint entrance animations, the engine handles initial state
# automatically — shapes referenced by entrance effects start hidden.
add_sequential_fade(s, groups)

set_notes(s, """ARCHITECTURE — ANIMATED (2m · Germaine):

This slide builds in three clicks so you can see the data path one stage at a time.

Click one — the CLIENT row appears. The user opens the site in any browser. They can record a 30-second video directly from their webcam using the built-in WebRTC APIs, or they can upload an existing file. Either way, the video is wrapped into a multipart HTTPS POST and sent to the backend.

Click two — CLOUDFLARE appears. The request hits Cloudflare's edge first. TLS terminates there. DDoS protection, WAF, and certificate management are all handled at the edge — I don't run any certificates on my homelab. From there, the request rides an outbound-only QUIC tunnel down to my house. There is no port forwarding. My home firewall stays closed. There is no public IP exposed.

Click three — the HOMELAB row appears. Three Docker containers under docker-compose. The cloudflared container terminates the tunnel and routes to nginx, which serves the React build. nginx reverse-proxies the /scan path to the FastAPI backend, which is where the rPPG pipeline runs. On the right, the same Python pipeline code is reused offline to produce the accuracy numbers we'll show on the next slide.

That's the entire production posture. Approximately three hundred lines of Python, one Dockerfile, one Cloudflare config. Nothing exotic.""")

# ============================================================
# SLIDE 6 — LIVE SITE / SCREENSHOTS (Germaine)
# ============================================================
s = add_slide()
header(s, "The live site",
       "vitalscan.bkre8tive.com  ·  open it on your phone right now",
       presenter="Germaine")

# four screenshot placeholders in 2x2 grid
def screenshot_box(left, top, w, h, caption):
    box = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, h)
    box.fill.solid(); box.fill.fore_color.rgb = SOFT_BG
    box.line.color.rgb = MUTED; box.line.width = Pt(1)
    add_text(s, left, top + h/2 - Inches(0.2), w, Inches(0.4),
             "[ drop screenshot here ]", size=11, color=MUTED,
             align=PP_ALIGN.CENTER)
    add_text(s, left, top + h - Inches(0.45), w, Inches(0.4),
             caption, size=11, bold=True, color=NAVY,
             align=PP_ALIGN.CENTER)

screenshot_box(Inches(0.5), Inches(1.6), Inches(6.0), Inches(2.5),
               "1 · Webcam capture (30 s countdown)")
screenshot_box(Inches(6.8), Inches(1.6), Inches(6.0), Inches(2.5),
               "2 · Body figure + biomarker cards")
screenshot_box(Inches(0.5), Inches(4.3), Inches(6.0), Inches(2.5),
               "3 · Heart rate, HRV, stress, BP")
screenshot_box(Inches(6.8), Inches(4.3), Inches(6.0), Inches(2.5),
               "4 · JSON response (developer view)")

set_notes(s, """LIVE SITE (90s · Germaine):

I want to step through what the site does in four images.

[Drop screenshots into the four boxes before presenting. To grab them: open vitalscan.bkre8tive.com in Chrome, Cmd-Shift-4 on Mac for region screenshot. Suggested four shots: (1) the landing page with the webcam button armed, (2) the dashboard with the body figure and biomarker cards, (3) zoom of the four biomarker cards showing real numbers, (4) the network tab showing the /scan POST and the JSON response.]

Image one — the landing page with the webcam recorder armed. The user clicks once. The browser asks for camera permission, then records for 30 seconds with a countdown.

Image two — the dashboard. After the scan finishes, the React app receives the biomarker JSON back from the backend and renders it as a body figure with marker positions and biomarker cards.

Image three — the four biomarker cards: heart rate in BPM, HRV in milliseconds, stress index between zero and one, and blood pressure. The first three come directly from the rPPG pipeline. Blood pressure is currently mocked — classical rPPG cannot derive blood pressure without a per-subject cuff calibration, so we return a plausible value with a known-limitation flag.

Image four — for developers, a view of the network tab showing the /scan POST, the multipart upload, and the JSON response. This is the shared API contract that any of the other groups can consume.

The site is live right now. If you have your phone out, open vitalscan.bkre8tive.com and try it.""")

# ============================================================
# SLIDE 7 — ACCURACY & PROGRESS (Abner / shared)
# ============================================================
s = add_slide()
header(s, "Accuracy and progress",
       "Where we are vs. where we want to be",
       presenter="Abner")

# headline number tile
tile = filled_card(s, Inches(0.5), Inches(1.6), Inches(6.0), Inches(2.6),
                   fill=GREEN_BG, border=GOOD)
add_text(s, Inches(0.75), Inches(1.8), Inches(5.5), Inches(0.4),
         "HEADLINE ACCURACY (Task 4)", size=11, bold=True, color=GOOD)
add_text(s, Inches(0.75), Inches(2.3), Inches(5.5), Inches(1.2),
         "4.06 BPM", size=60, bold=True, color=NAVY)
add_text(s, Inches(0.75), Inches(3.4), Inches(5.5), Inches(0.4),
         "Mean abs. error on 42 UBFC subjects",
         size=14, color=SLATE)
add_text(s, Inches(0.75), Inches(3.8), Inches(5.5), Inches(0.4),
         "Rubric target < 10 BPM   ·   median 1.05 BPM",
         size=12, color=MUTED)

# progress checklist
filled_card(s, Inches(6.8), Inches(1.6), Inches(6.0), Inches(5.0),
            fill=SOFT_BG, border=ACCENT)
add_text(s, Inches(7.05), Inches(1.8), Inches(5.5), Inches(0.4),
         "PROGRESS CHECKLIST", size=11, bold=True, color=ACCENT)

def progress_row(top, mark, color, text):
    sym = s.shapes.add_shape(MSO_SHAPE.OVAL,
                            Inches(7.1), top, Inches(0.3), Inches(0.3))
    sym.fill.solid(); sym.fill.fore_color.rgb = color
    sym.line.fill.background()
    add_text(s, Inches(7.1), top, Inches(0.3), Inches(0.3),
             mark, size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, Inches(7.55), top + Inches(0.02), Inches(5.0), Inches(0.4),
             text, size=12, color=INK)

progress_row(Inches(2.4), "✓", GOOD, "Datasets retrieved (UBFC + SCAMPS)")
progress_row(Inches(2.9), "✓", GOOD, "rPPG pipeline (face → POS → BPM)")
progress_row(Inches(3.4), "✓", GOOD, "FastAPI /scan endpoint live")
progress_row(Inches(3.9), "✓", GOOD, "React dashboard + webcam capture")
progress_row(Inches(4.4), "✓", GOOD, "Deployed to vitalscan.bkre8tive.com")
progress_row(Inches(4.9), "✓", GOOD, "Accuracy report (MAE / RMSE)")
progress_row(Inches(5.4), "•", WARN, "Blood pressure — still mocked")
progress_row(Inches(5.9), "•", WARN, "Deep-learning stretch — scaffolded only")

# bottom line
add_text(s, Inches(0.5), Inches(4.6), Inches(6.0), Inches(0.4),
         "Distribution on UBFC:", size=12, bold=True, color=NAVY)
add_bullets(s, Inches(0.5), Inches(5.0), Inches(6.0), Inches(1.8), [
    "50% of subjects within ± 1 BPM",
    "67% within ± 3 BPM",
    "90% within ± 10 BPM",
], size=12)

set_notes(s, """ACCURACY & PROGRESS (90s · Abner or shared):

The big green number on the left is the headline. On the 42-subject UBFC-rPPG dataset — that's the real-human dataset with synchronized pulse-oximeter ground truth — our pipeline measures heart rate with a mean absolute error of 4.06 beats per minute. The course rubric asked for under 10. We're at roughly forty percent of that target. The median error is just over one beat per minute, meaning half of our subjects are read essentially perfectly.

The distribution at the bottom of the left column is what we'd cite if someone asks whether the average is hiding bad outliers. Half the subjects within one BPM. Two-thirds within three. Ninety percent within ten. The errors are not in a long tail — most subjects are read accurately, with a small cluster of four subjects where both algorithms struggle. That cluster correlates with a recording-session issue in the dataset itself.

The right column is the project status. Six things complete. Two known limits. Blood pressure is still mocked — that's a real constraint of classical rPPG. And the deep-learning extension was a stretch goal in the brief; the scaffolding is in our repo but we haven't trained the model.

We're treating this as a working-progress presentation. The deliverable is done. What we're going to refine before final submission is the writeup polish and maybe one more pass at the outliers.""")

# ============================================================
# SLIDE 8 — THANK YOU / Q&A
# ============================================================
s = add_slide()
panel = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
panel.line.fill.background()
panel.fill.solid(); panel.fill.fore_color.rgb = NAVY

add_text(s, Inches(0.6), Inches(0.8), Inches(12), Inches(0.5),
         "AIT 500 · WESTCLIFF UNIVERSITY · GROUP 1", size=12, bold=True, color=ACCENT)
add_text(s, Inches(0.6), Inches(1.6), Inches(12), Inches(1.3),
         "Thank you — questions?", size=58, bold=True, color=WHITE)

add_text(s, Inches(0.6), Inches(3.4), Inches(12), Inches(0.5),
         "Try the site right now", size=16, bold=True, color=ACCENT)
add_text(s, Inches(0.6), Inches(3.9), Inches(12), Inches(0.6),
         "vitalscan.bkre8tive.com",
         size=34, bold=True, color=WHITE)

add_text(s, Inches(0.6), Inches(5.0), Inches(12), Inches(0.5),
         "Code · github.com/gbiz21/vitalscan",
         size=14, color=ACCENT)
add_text(s, Inches(0.6), Inches(5.55), Inches(12), Inches(0.5),
         "Group 1 — Germaine · Jason · Daray · Abner",
         size=14, color=WHITE)

set_notes(s, """CLOSE (30s):

Recap in two sentences. We took the Group 1 assignment — rPPG signal extraction — and built it out into a full working product: pipeline, REST API, dashboard, deployed to a live URL. On the 42-subject UBFC real-human dataset we measure 4.06 BPM mean absolute error, well inside the course target.

The site is live at vitalscan.bkre8tive.com. The code is on GitHub. Thank you. Happy to take questions.""")


# write
prs.save(OUT)
print(f"wrote {OUT}")
