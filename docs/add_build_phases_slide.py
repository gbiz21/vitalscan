"""Add a 'Build phases' slide to the backend-architecture deck.

Inserts a new slide between the existing slide 6 (The live site) and
slide 7 (Accuracy and progress) so the narrative becomes:
    ... → architecture → live site → BUILD PHASES → accuracy → thank you

The 5 phases are taken verbatim from the project README's Build Phases table.
"""
from copy import deepcopy
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

DECK = Path(__file__).resolve().parent / "VitalScan_Group1_Backend_System_Architecture.pptx"

# Palette — match the rest of the deck.
NAVY = RGBColor(0x0F, 0x1F, 0x44)
INK = RGBColor(0x1A, 0x1A, 0x1A)
MUTED = RGBColor(0x55, 0x66, 0x7A)
ACCENT = RGBColor(0x1F, 0x6F, 0xEB)
GOOD = RGBColor(0x1E, 0x88, 0x4F)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SOFT_BG = RGBColor(0xF6, 0xF8, 0xFB)
CARD_BG = RGBColor(0xEE, 0xF3, 0xFA)

PHASES = [
    ("Phase 1", "React UI with mock biomarker data",
     "Got the dashboard, body figure, and biomarker cards working against fake data first."),
    ("Phase 2", "FastAPI backend with mock /scan",
     "Stood up the /scan endpoint returning the shared JSON contract — no rPPG yet."),
    ("Phase 3", "Real rPPG pipeline",
     "Replaced the mock with the actual face detection · POS · HRV signal chain."),
    ("Phase 4", "Docker + Cloudflare Tunnel deployment",
     "Three-container stack on the homelab, live at vitalscan.bkre8tive.com."),
    ("Phase 5", "Browser webcam capture",
     "In-browser 30-second capture using getUserMedia and MediaRecorder."),
]


def add_text(slide, left, top, width, height, text, *, size=18, bold=False,
             color=INK, align=PP_ALIGN.LEFT, font="Calibri"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.name = font; r.font.size = Pt(size)
        r.font.bold = bold; r.font.color.rgb = color
    return tb


def add_header(slide, title, subtitle, presenter):
    stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0,
                                    Inches(13.333), Inches(0.18))
    stripe.line.fill.background()
    stripe.fill.solid(); stripe.fill.fore_color.rgb = ACCENT
    add_text(slide, Inches(0.5), Inches(0.36), Inches(12.3), Inches(0.45),
             title, size=28, bold=True, color=NAVY)
    add_text(slide, Inches(0.58), Inches(0.89), Inches(8.5), Inches(0.3),
             subtitle, size=14, color=MUTED)
    add_text(slide, Inches(9.35), Inches(0.42), Inches(2.45), Inches(0.25),
             "VitalScan · AIT 500 · Group 1 · Week 8", size=10, color=MUTED)
    add_text(slide, Inches(11.85), Inches(0.42), Inches(0.95), Inches(0.25),
             f"presented by {presenter}", size=10, color=ACCENT,
             align=PP_ALIGN.RIGHT)
    add_text(slide, Inches(0.5), Inches(7.22), Inches(8.0), Inches(0.2),
             "vitalscan.bkre8tive.com  ·  github.com/gbiz21/vitalscan",
             size=10, color=MUTED)


def add_phase_card(slide, top, phase, headline, detail):
    # Card background
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  Inches(0.75), top,
                                  Inches(11.8), Inches(0.85))
    card.adjustments[0] = 0.18
    card.fill.solid(); card.fill.fore_color.rgb = SOFT_BG
    card.line.color.rgb = ACCENT; card.line.width = Pt(1)

    # Done marker — green circle with check
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                    Inches(1.0), top + Inches(0.21),
                                    Inches(0.42), Inches(0.42))
    circle.fill.solid(); circle.fill.fore_color.rgb = GOOD
    circle.line.fill.background()
    add_text(slide, Inches(1.0), top + Inches(0.22), Inches(0.42), Inches(0.42),
             "✓", size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Phase label
    add_text(slide, Inches(1.65), top + Inches(0.18), Inches(1.3), Inches(0.4),
             phase, size=14, bold=True, color=ACCENT)
    # Headline
    add_text(slide, Inches(3.0), top + Inches(0.16), Inches(9.4), Inches(0.4),
             headline, size=16, bold=True, color=NAVY)
    # Detail
    add_text(slide, Inches(3.0), top + Inches(0.5), Inches(9.4), Inches(0.3),
             detail, size=11, color=MUTED)


def build_slide(prs):
    # This deck only has one layout — index 0 (the default blank layout).
    BLANK = prs.slide_layouts[0]
    s = prs.slides.add_slide(BLANK)

    # White background
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0,
                            prs.slide_width, prs.slide_height)
    bg.line.fill.background()
    bg.fill.solid(); bg.fill.fore_color.rgb = WHITE

    add_header(s, "Build phases",
               "How we got from project brief to a live, working product",
               presenter="Germaine")

    # 5 phase cards, evenly spaced from T=1.4 to T=6.8
    top = Inches(1.4)
    for phase, headline, detail in PHASES:
        add_phase_card(s, top, phase, headline, detail)
        top += Inches(1.05)

    # Speaker notes — this deck's notes master doesn't always expose a notes
    # placeholder; wrap so a missing placeholder doesn't kill the build.
    notes_body = (
        "BUILD PHASES (90s · Germaine):\n\n"
        "We didn't try to build everything at once — we did it in five phases "
        "and validated each before moving to the next.\n\n"
        "Phase 1 — React UI with mock biomarker data. We got the dashboard, "
        "body figure, and biomarker cards rendering against hard-coded JSON. "
        "This let us nail the user experience before any of the heavy lifting.\n\n"
        "Phase 2 — FastAPI backend with a mock /scan endpoint. Same JSON "
        "contract, returned from a real HTTP service. The frontend swapped "
        "from local mock data to the API with no UI changes — that's the value "
        "of agreeing on a contract first.\n\n"
        "Phase 3 — the real rPPG pipeline. Face detection, ROI extraction, "
        "POS algorithm, HRV — all the stuff Daray just walked through. The "
        "endpoint signature didn't change; only the implementation behind it.\n\n"
        "Phase 4 — Dockerized everything and deployed to the homelab through a "
        "Cloudflare Tunnel. That's how vitalscan.bkre8tive.com became reachable.\n\n"
        "Phase 5 — browser webcam capture. The last piece — letting people "
        "scan themselves on the live site without uploading a file.\n\n"
        "All five phases are complete. The site you can hit right now is "
        "running phase 5 against phase 3."
    )
    try:
        notes_tf = s.notes_slide.notes_text_frame
        if notes_tf is not None:
            notes_tf.text = notes_body
    except (AttributeError, ValueError):
        # Notes master lacks a placeholder — skip; on-slide content stands alone.
        pass
    return s


def move_slide_to_position(prs, slide_idx_from, slide_idx_to):
    """Reorder slides via the sldIdLst element."""
    sldIdLst = prs.slides._sldIdLst
    children = list(sldIdLst)
    moved = children[slide_idx_from]
    sldIdLst.remove(moved)
    sldIdLst.insert(slide_idx_to, moved)


def main():
    prs = Presentation(DECK)
    n_before = len(prs.slides)
    print(f"slides before: {n_before}")

    # Skip if already present (idempotency check)
    for i, s in enumerate(prs.slides):
        for sh in s.shapes:
            if sh.has_text_frame and sh.text_frame.text.strip() == "Build phases":
                print(f"  build-phases slide already present (slide {i+1}) — replacing.")
                # Drop existing slide; easiest by removing its sldId entry.
                # Note: this leaves the underlying part orphaned, but PowerPoint
                # is tolerant of that and we keep the deck clean enough.
                sldIdLst = prs.slides._sldIdLst
                sldIdLst.remove(list(sldIdLst)[i])
                break

    build_slide(prs)
    # New slide is appended at the end. Move it to position 6 (0-indexed),
    # i.e. right after "The live site" and before "Accuracy and progress".
    new_index = len(prs.slides) - 1
    target_index = 6  # 0-indexed: after slide 6 (which is index 5)
    move_slide_to_position(prs, new_index, target_index)

    prs.save(DECK)
    n_after = len(prs.slides)
    print(f"slides after:  {n_after}")
    print(f"  inserted 'Build phases' at position 7 (between 'Live site' and 'Accuracy')")


if __name__ == "__main__":
    main()
