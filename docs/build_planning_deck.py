"""Build the THIS WEEK planning & architecture blueprint deck.

This is distinct from the Week 8 results deck:
  - Week 8 deck (Backend_System_Architecture.pptx): shows what was built.
  - This deck (Planning_Architecture.pptx): shows the system DESIGN
    and BUILD ROADMAP — the blueprint the professor specifically asked
    for ("focus on designing out the system... architecting the system
    itself" from the 2026-05-13 class).

Run:  backend/venv312/bin/python docs/build_planning_deck.py
"""
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"

ROOT = Path(__file__).resolve().parent.parent
FIG_DIR = ROOT / "docs" / "figures"
OUT = ROOT / "docs" / "VitalScan_Group1_Planning_Architecture.pptx"

# --- palette (same as the other decks for visual consistency) ---
NAVY = RGBColor(0x0F, 0x1F, 0x44)
SLATE = RGBColor(0x33, 0x44, 0x5C)
INK = RGBColor(0x1A, 0x1A, 0x1A)
MUTED = RGBColor(0x55, 0x66, 0x7A)
ACCENT = RGBColor(0x1F, 0x6F, 0xEB)
GOOD = RGBColor(0x1E, 0x88, 0x4F)
WARN = RGBColor(0xC0, 0x73, 0x1D)
DANGER = RGBColor(0xC0, 0x39, 0x2B)
SOFT_BG = RGBColor(0xF6, 0xF8, 0xFB)
CARD_BG = RGBColor(0xEE, 0xF3, 0xFA)
GREEN_BG = RGBColor(0xE9, 0xF5, 0xEC)
ORANGE_BG = RGBColor(0xFD, 0xF1, 0xE3)
RED_BG = RGBColor(0xFA, 0xEC, 0xEA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


# ---------- helpers ----------
def add_slide():
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE
    return s


def add_text(slide, left, top, width, height, text, *, size=18, bold=False,
             color=INK, align=PP_ALIGN.LEFT, font="Calibri"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.name = font; r.font.size = Pt(size)
        r.font.bold = bold; r.font.color.rgb = color
    return tb


def add_bullets(slide, left, top, width, height, bullets, *, size=14, color=INK):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(6)
        r = p.add_run()
        r.text = f"•  {b}"
        r.font.name = "Calibri"
        r.font.size = Pt(size)
        r.font.color.rgb = color
    return tb


def filled_card(slide, left, top, width, height, *,
                fill=SOFT_BG, border=ACCENT, radius=0.06):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  left, top, width, height)
    card.adjustments[0] = radius
    card.fill.solid()
    card.fill.fore_color.rgb = fill
    card.line.color.rgb = border
    card.line.width = Pt(1.25)
    card.shadow.inherit = False
    return card


def header(slide, title, subtitle=None):
    stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, Inches(0.18))
    stripe.line.fill.background()
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = ACCENT
    add_text(slide, Inches(0.5), Inches(0.35), Inches(12.3), Inches(0.7),
             title, size=28, bold=True, color=NAVY)
    if subtitle:
        add_text(slide, Inches(0.5), Inches(0.95), Inches(12.3), Inches(0.4),
                 subtitle, size=14, color=MUTED)
    add_text(slide, Inches(0.5), Inches(7.05), Inches(12.3), Inches(0.3),
             "VitalScan · AIT 500 · Group 1 · Planning checkpoint",
             size=10, color=MUTED)


def set_notes(slide, text):
    try:
        notes_tf = slide.notes_slide.notes_text_frame
        if notes_tf is not None:
            notes_tf.text = text
    except (AttributeError, ValueError):
        pass


# ---------- animation injection (sequential fade-in on click) ----------
def add_sequential_fade(slide, groups, *, dur_ms=600):
    """For each inner list of shape ids, fade those shapes in on the next click."""
    p = "{%s}" % P_NS
    timing = etree.SubElement(slide._element, f"{p}timing")
    tnLst = etree.SubElement(timing, f"{p}tnLst")
    root_par = etree.SubElement(tnLst, f"{p}par")
    root_ctn = etree.SubElement(
        root_par, f"{p}cTn",
        id="1", dur="indefinite", restart="never", nodeType="tmRoot",
    )
    root_children = etree.SubElement(root_ctn, f"{p}childTnLst")
    seq = etree.SubElement(root_children, f"{p}seq",
                           concurrent="1", nextAc="seek")
    main_ctn = etree.SubElement(seq, f"{p}cTn",
                                id="2", dur="indefinite", nodeType="mainSeq")
    main_children = etree.SubElement(main_ctn, f"{p}childTnLst")
    prev_cond = etree.SubElement(seq, f"{p}prevCondLst")
    cond1 = etree.SubElement(prev_cond, f"{p}cond", evt="onPrev", delay="0")
    etree.SubElement(etree.SubElement(cond1, f"{p}tgtEl"), f"{p}sldTgt")
    next_cond = etree.SubElement(seq, f"{p}nextCondLst")
    cond2 = etree.SubElement(next_cond, f"{p}cond", evt="onNext", delay="0")
    etree.SubElement(etree.SubElement(cond2, f"{p}tgtEl"), f"{p}sldTgt")

    next_id = [3]
    def nid():
        i = next_id[0]; next_id[0] += 1
        return str(i)

    for shape_ids in groups:
        click_par = etree.SubElement(main_children, f"{p}par")
        click_ctn = etree.SubElement(click_par, f"{p}cTn", id=nid(), fill="hold")
        etree.SubElement(etree.SubElement(click_ctn, f"{p}stCondLst"),
                         f"{p}cond", delay="indefinite")
        click_children = etree.SubElement(click_ctn, f"{p}childTnLst")
        grp_par = etree.SubElement(click_children, f"{p}par")
        grp_ctn = etree.SubElement(grp_par, f"{p}cTn", id=nid(), fill="hold")
        etree.SubElement(etree.SubElement(grp_ctn, f"{p}stCondLst"),
                         f"{p}cond", delay="0")
        grp_children = etree.SubElement(grp_ctn, f"{p}childTnLst")
        for idx, sid in enumerate(shape_ids):
            node_type = "clickEffect" if idx == 0 else "withEffect"
            effect_par = etree.SubElement(grp_children, f"{p}par")
            effect_ctn = etree.SubElement(
                effect_par, f"{p}cTn",
                id=nid(), presetID="10", presetClass="entr",
                presetSubtype="0", fill="hold", grpId="0",
                nodeType=node_type,
            )
            etree.SubElement(etree.SubElement(effect_ctn, f"{p}stCondLst"),
                             f"{p}cond", delay="0")
            effect_children = etree.SubElement(effect_ctn, f"{p}childTnLst")
            set_el = etree.SubElement(effect_children, f"{p}set")
            cbhvr = etree.SubElement(set_el, f"{p}cBhvr")
            etree.SubElement(cbhvr, f"{p}cTn", id=nid(), dur="1", fill="hold")
            tgt = etree.SubElement(cbhvr, f"{p}tgtEl")
            etree.SubElement(tgt, f"{p}spTgt", spid=str(sid))
            an = etree.SubElement(etree.SubElement(cbhvr, f"{p}attrNameLst"),
                                  f"{p}attrName")
            an.text = "style.visibility"
            etree.SubElement(etree.SubElement(set_el, f"{p}to"),
                             f"{p}strVal", val="visible")
            anim = etree.SubElement(effect_children, f"{p}anim",
                                    calcmode="lin", valueType="num")
            cbhvr2 = etree.SubElement(anim, f"{p}cBhvr", additive="base")
            etree.SubElement(cbhvr2, f"{p}cTn",
                             id=nid(), dur=str(dur_ms), fill="hold")
            tgt2 = etree.SubElement(cbhvr2, f"{p}tgtEl")
            etree.SubElement(tgt2, f"{p}spTgt", spid=str(sid))
            an2 = etree.SubElement(etree.SubElement(cbhvr2, f"{p}attrNameLst"),
                                   f"{p}attrName")
            an2.text = "style.opacity"
            tavLst = etree.SubElement(anim, f"{p}tavLst")
            for tm, val in [("0", "0"), ("100000", "1")]:
                tav = etree.SubElement(tavLst, f"{p}tav", tm=tm)
                etree.SubElement(etree.SubElement(tav, f"{p}val"),
                                 f"{p}fltVal", val=val)


# ============================================================
# SLIDE 1 — TITLE
# ============================================================
s = add_slide()
panel = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(5.6), SH)
panel.line.fill.background()
panel.fill.solid()
panel.fill.fore_color.rgb = NAVY

add_text(s, Inches(0.6), Inches(0.6), Inches(4.8), Inches(0.4),
         "AIT 500 · WESTCLIFF UNIVERSITY", size=11, bold=True, color=ACCENT)
add_text(s, Inches(0.6), Inches(1.3), Inches(5.0), Inches(1.6),
         "VitalScan", size=66, bold=True, color=WHITE)
add_text(s, Inches(0.6), Inches(2.8), Inches(5.0), Inches(0.8),
         "Planning &\nArchitecture Blueprint",
         size=22, color=WHITE)
add_text(s, Inches(0.6), Inches(4.6), Inches(4.8), Inches(0.4),
         "GROUP 1 · rPPG SIGNAL EXTRACTION",
         size=12, bold=True, color=ACCENT)
add_text(s, Inches(0.6), Inches(6.6), Inches(4.8), Inches(0.4),
         "Class planning checkpoint", size=13, bold=True, color=ACCENT)

# right side — what's in this deck
add_text(s, Inches(6.2), Inches(0.9), Inches(6.5), Inches(0.4),
         "WHAT THIS DECK COVERS", size=12, bold=True, color=MUTED)
add_text(s, Inches(6.2), Inches(1.35), Inches(6.5), Inches(0.7),
         "Designing before coding",
         size=30, bold=True, color=NAVY)

def design_card(top, label, body):
    card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                              Inches(6.2), top, Inches(6.3), Inches(0.85))
    card.adjustments[0] = 0.18
    card.fill.solid(); card.fill.fore_color.rgb = CARD_BG
    card.line.color.rgb = ACCENT; card.line.width = Pt(1)
    add_text(s, Inches(6.45), top + Inches(0.12),
             Inches(2.6), Inches(0.4), label, size=14, bold=True, color=NAVY)
    add_text(s, Inches(6.45), top + Inches(0.45),
             Inches(5.7), Inches(0.4), body, size=11, color=SLATE)

design_card(Inches(2.55), "1 · The problem",      "What Group 1 must deliver and why")
design_card(Inches(3.50), "2 · System architecture","Components, layers, data flow")
design_card(Inches(4.45), "3 · Build roadmap",    "Phase 1 – 5, planned in order")
design_card(Inches(5.40), "4 · Risks & next steps","What could go wrong, what we do next")

add_text(s, Inches(6.2), Inches(6.55), Inches(6.5), Inches(0.4),
         "Germaine · Jason · Daray · Abner   ·   AIT 500 · 2026",
         size=12, color=MUTED)

set_notes(s, """OPENER (45s):

Good [morning/afternoon]. We are Group 1 for the VitalScan project — Germaine, Jason, Daray, and Abner. Our slice of the project is rPPG signal extraction.

This deck is intentionally a planning checkpoint, not a results presentation. The professor asked us to focus on systems engineering rather than jumping straight into code — design the system first, then implement.

So in the next ten minutes we'll walk through four things: the problem we're solving, the system architecture we've designed, the build phases we'll execute in order, and the risks we anticipate. We'll save the actual results for the Week 8 group presentation.""")

# ============================================================
# SLIDE 2 — THE PROBLEM / OBJECTIVES
# ============================================================
s = add_slide()
header(s, "The problem we're solving",
       "Group 1's assignment from the project brief — what we must deliver")

# What goes in / out
filled_card(s, Inches(0.5), Inches(1.7), Inches(6.0), Inches(2.5))
add_text(s, Inches(0.75), Inches(1.85), Inches(5.5), Inches(0.4),
         "INPUT", size=12, bold=True, color=MUTED)
add_text(s, Inches(0.75), Inches(2.25), Inches(5.5), Inches(0.5),
         "A 30-second facial video", size=20, bold=True, color=NAVY)
add_bullets(s, Inches(0.75), Inches(2.85), Inches(5.5), Inches(1.3), [
    "Standard webcam · phone or laptop",
    "30 fps · no special hardware",
    "User in front of the camera, well-lit",
])

filled_card(s, Inches(6.8), Inches(1.7), Inches(6.0), Inches(2.5),
            fill=GREEN_BG, border=GOOD)
add_text(s, Inches(7.05), Inches(1.85), Inches(5.5), Inches(0.4),
         "OUTPUT", size=12, bold=True, color=GOOD)
add_text(s, Inches(7.05), Inches(2.25), Inches(5.5), Inches(0.5),
         "Biomarker JSON contract", size=20, bold=True, color=NAVY)
add_bullets(s, Inches(7.05), Inches(2.85), Inches(5.5), Inches(1.3), [
    "Heart rate (BPM) · HRV (SDNN, ms)",
    "Stress index (0 – 1)",
    "Blood pressure (placeholder, flagged)",
])

# Objectives
add_text(s, Inches(0.5), Inches(4.4), Inches(12.3), Inches(0.4),
         "OBJECTIVES (per the project brief)",
         size=12, bold=True, color=MUTED)

def obj_card(left, num, title, body):
    card = filled_card(s, left, Inches(4.85), Inches(2.95), Inches(1.85))
    add_text(s, left + Inches(0.2), Inches(5.0), Inches(0.6), Inches(0.45),
             num, size=24, bold=True, color=ACCENT)
    add_text(s, left + Inches(0.2), Inches(5.5), Inches(2.55), Inches(0.4),
             title, size=14, bold=True, color=NAVY)
    add_text(s, left + Inches(0.2), Inches(5.85), Inches(2.55), Inches(0.85),
             body, size=11, color=SLATE)

obj_card(Inches(0.5),  "1",
         "Face detection", "Detect the user's face on every video frame.")
obj_card(Inches(3.65), "2",
         "ROI extraction", "Pick skin patches — forehead + cheeks — with the best signal.")
obj_card(Inches(6.80), "3",
         "rPPG signal",    "Recover heart rate, HRV, and stress from the ROI color.")
obj_card(Inches(9.95), "4",
         "API contract",   "Expose results as JSON so other groups can consume them.")

set_notes(s, """THE PROBLEM (90s):

Two-column framing first — input and output. Our input is a 30-second video from any standard webcam, no special hardware required. Our output is a biomarker JSON object: heart rate in BPM, heart-rate variability in milliseconds, a stress index from zero to one, and blood pressure. Blood pressure we'll flag as a placeholder because classical rPPG can't actually derive it — more on that in the risks slide.

The four numbered objectives below come straight from the project brief. Group 1's slice is the rPPG signal extraction — face detection, region-of-interest picking, the signal-recovery algorithm, and exposing the result as a JSON contract that the other groups can integrate against.

The contract is the part worth emphasizing. The professor said "be systems engineers, not just programmers" — agreeing on the API shape early is the most concrete way to design before coding. Everything downstream of that contract can be built in isolation.""")

# ============================================================
# SLIDE 3 — SYSTEM ARCHITECTURE (use figure 2)
# ============================================================
s = add_slide()
header(s, "System architecture — the blueprint",
       "Three layers, one request path, fully containerized")

s.shapes.add_picture(str(FIG_DIR / "figure2-architecture.png"),
                     Inches(0.5), Inches(1.55), width=Inches(12.3))

add_text(s, Inches(0.5), Inches(6.55), Inches(12.3), Inches(0.4),
         "Client (browser) → Cloudflare edge + tunnel → homelab docker-compose stack",
         size=13, bold=True, color=MUTED, align=PP_ALIGN.CENTER)

set_notes(s, """SYSTEM ARCHITECTURE (2m):

This is our system design end-to-end. Walking left-to-right, top-to-bottom.

Top — the client. Any device with a browser. Three input modes we're planning: in-browser webcam capture using getUserMedia and MediaRecorder, file upload for users who recorded elsewhere, and the dashboard view that displays results. All in standard React running on the user's machine.

Middle — Cloudflare. The professor emphasized data quality and security; this is how we handle TLS, DDoS protection, and edge routing without managing certificates ourselves. Cloudflare's named tunnel is outbound-only — our homelab firewall stays closed, no public IP exposed. This is a deliberate security-first design choice.

Bottom — the homelab host. Three containers under docker-compose. The cloudflared container terminates the tunnel and routes to nginx, which serves the React build and reverse-proxies the /scan endpoint to the FastAPI backend. The backend is where the actual rPPG pipeline runs — face detection, POS algorithm, HRV calculation.

Right side — the offline evaluation harness. Same Python code path as the runtime, so accuracy numbers we report are produced by the same code that serves real traffic. That's intentional — no "demo code vs production code" gap.

Design rationale to emphasize: separation of concerns. The frontend doesn't know how rPPG works. The backend doesn't know what UI exists. The deployment doesn't know the algorithm. Each layer can be swapped independently.""")

# ============================================================
# SLIDE 4 — rPPG SIGNAL CHAIN  [ANIMATED · 4 click reveals]
# ============================================================
s = add_slide()
header(s, "rPPG signal pipeline — algorithm blueprint",
       "Click to step through · four stages from video to biomarkers")

groups_slide4 = []  # one inner list per click reveal

def numbered_step_card(left, num, title, body, *, accent=ACCENT):
    """Return list of shape ids for the card + its arrow gets added later."""
    ids = []
    card = filled_card(s, left, Inches(1.65), Inches(2.9), Inches(4.85),
                       fill=SOFT_BG, border=accent)
    ids.append(card.shape_id)
    circle = s.shapes.add_shape(MSO_SHAPE.OVAL,
                               left + Inches(0.25), Inches(1.85),
                               Inches(0.7), Inches(0.7))
    circle.fill.solid(); circle.fill.fore_color.rgb = accent
    circle.line.fill.background()
    ids.append(circle.shape_id)
    num_tb = add_text(s, left + Inches(0.25), Inches(1.9),
                     Inches(0.7), Inches(0.6),
                     str(num), size=22, bold=True, color=WHITE,
                     align=PP_ALIGN.CENTER)
    ids.append(num_tb.shape_id)
    title_tb = add_text(s, left + Inches(0.25), Inches(2.75),
                        Inches(2.4), Inches(0.4),
                        title, size=15, bold=True, color=NAVY)
    ids.append(title_tb.shape_id)
    body_tb = add_bullets(s, left + Inches(0.25), Inches(3.45),
                          Inches(2.4), Inches(3.0),
                          body, size=11)
    ids.append(body_tb.shape_id)
    return ids


step1 = numbered_step_card(Inches(0.5), 1, "Detect face", [
    "MediaPipe FaceMesh",
    "468 landmarks / frame",
    "Real-time on CPU",
    "Skip frames with no face",
])
groups_slide4.append(step1)

step2 = numbered_step_card(Inches(3.6), 2, "Pick ROIs", [
    "Forehead + both cheeks",
    "Highest capillary density",
    "Lowest motion artifact",
    "Mean RGB per frame → T×3",
])
groups_slide4.append(step2)

step3 = numbered_step_card(Inches(6.7), 3, "POS + bandpass", [
    "Project RGB onto plane",
    "orthogonal to skin tone",
    "Butterworth 0.7–4 Hz",
    "Pulse amplified, noise out",
])
groups_slide4.append(step3)

step4 = numbered_step_card(Inches(9.8), 4, "FFT → BPM", [
    "FFT with 4× zero-pad",
    "Search peak in 60–210 BPM",
    "Peak × 60 = heart rate",
    "Reuse pulse for HRV + stress",
], accent=GOOD)
groups_slide4.append(step4)

# arrows BETWEEN steps — each arrow reveals with the right-hand step
for arrow_left, group_idx in [(Inches(3.35), 1), (Inches(6.45), 2), (Inches(9.55), 3)]:
    arr = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                             arrow_left, Inches(3.85),
                             Inches(0.3), Inches(0.4))
    arr.fill.solid(); arr.fill.fore_color.rgb = MUTED
    arr.line.fill.background()
    groups_slide4[group_idx].append(arr.shape_id)

add_text(s, Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.4),
         "Same cleaned pulse → HR · HRV (SDNN) · stress index (LF/HF sigmoid)",
         size=12, color=MUTED, align=PP_ALIGN.CENTER)

add_sequential_fade(s, groups_slide4)

set_notes(s, """rPPG SIGNAL PIPELINE (2m):

This is the algorithm we'll implement, designed before any code is written.

The physics: a heartbeat pushes oxygenated blood through facial capillaries. Hemoglobin absorbs green light. So the amount of green reflected from skin dips slightly with each beat — about 0.1 to 1 percent of total reflected light. Invisible to the eye, recoverable from a webcam if we process enough pixels over enough frames.

Stage 1 — face detection and ROI. We'll use MediaPipe FaceMesh, which puts 468 landmarks on a face in real time on CPU. We pick forehead and both cheeks — highest capillary density, lowest motion. Output is a time-series of average RGB values.

Stage 2 — POS algorithm from Wang et al. 2017. The trick: lighting noise moves all three RGB channels together; the pulse signal lives mostly in green. So we project the data onto the plane orthogonal to the skin-tone vector. Noise drops out, pulse is amplified. Then bandpass filter at 0.7 to 4 Hz, take an FFT, find the dominant peak between 60 and 210 BPM.

Stage 3 — HRV and stress. Peak-detect the cleaned waveform to get inter-beat intervals. Standard deviation of those intervals is SDNN. The low-frequency over high-frequency power ratio gives us the stress index.

Stage 4 — wrap it in JSON matching the shared contract.

This is the design. The validation strategy — covered on the data plan slide — is to run this against UBFC-rPPG with ground-truth pulse oximetry and confirm MAE under 10 BPM per the rubric.""")

# ============================================================
# SLIDE 5 — TECH STACK & DESIGN CHOICES
# ============================================================
s = add_slide()
header(s, "Tech stack & design choices",
       "Each tool picked for a reason — kept boring and replaceable")

cols = [
    ("FRONTEND", ACCENT, [
        ("React 18 + TypeScript", "Component model + type safety"),
        ("Vite",                  "Fast dev server + production build"),
        ("Tailwind CSS",          "Utility classes — no design bikeshed"),
    ]),
    ("BACKEND",  GOOD, [
        ("Python 3.12 + FastAPI", "Async REST in <50 lines"),
        ("numpy + scipy",         "Vectorized signal-chain math"),
        ("MediaPipe FaceMesh",    "Real-time landmarks on CPU"),
    ]),
    ("DEPLOY",   WARN, [
        ("Docker Compose",        "Three services, one command"),
        ("Cloudflare Tunnel",     "TLS + no port forwarding"),
        ("Homelab host",          "Free, fast iteration"),
    ]),
]

col_w = Inches(4.0)
gap = Inches(0.2)
left = Inches(0.5)

for col_idx, (title, color, items) in enumerate(cols):
    x = left + (col_w + gap) * col_idx
    filled_card(s, x, Inches(1.7), col_w, Inches(5.0),
                fill=SOFT_BG, border=color)
    add_text(s, x + Inches(0.25), Inches(1.85), col_w - Inches(0.5), Inches(0.4),
             title, size=14, bold=True, color=color)

    y = Inches(2.4)
    for tool, why in items:
        add_text(s, x + Inches(0.25), y, col_w - Inches(0.5), Inches(0.35),
                 tool, size=14, bold=True, color=NAVY)
        add_text(s, x + Inches(0.25), y + Inches(0.35), col_w - Inches(0.5), Inches(0.6),
                 why, size=11, color=SLATE)
        y += Inches(1.05)

set_notes(s, """TECH STACK (90s):

Three columns — frontend, backend, deployment.

Frontend: React with TypeScript. We could have used vanilla JS, but typed component models scale better when multiple people touch the code. Vite for builds because it's the modern standard and the dev server is instant. Tailwind for styling because it eliminates the "what should this button look like" conversation — utility classes encode the design system.

Backend: Python 3.12 because that's what the rPPG and ML ecosystem speaks. FastAPI for the HTTP layer because it gives us async, automatic OpenAPI docs, and a 50-line endpoint. Numpy and scipy for the actual signal-chain math — vectorized, fast, well-tested. MediaPipe for face detection because it's the only library that runs at video framerate on CPU.

Deployment: Docker Compose because three containers is the right granularity — frontend, backend, tunnel. Cloudflare Tunnel because it solves three problems at once: TLS termination, DDoS protection, and no port forwarding. Homelab because free and we iterate faster on hardware we own.

Design principle we're applying throughout: keep it boring and keep it replaceable. No fancy frameworks, no novel tools. If any piece needs to swap out, it should be a one-day job.""")

# ============================================================
# SLIDE 6 — BUILD PHASES ROADMAP
# ============================================================
s = add_slide()
header(s, "Build phases — roadmap",
       "Five phases, validated at each boundary before moving on")

phases = [
    ("Phase 1", "React UI with mock biomarker data",
     "Nail the dashboard + body figure + biomarker cards against hard-coded JSON. "
     "De-risk the UX before any of the heavy lifting."),
    ("Phase 2", "FastAPI backend with mock /scan",
     "Stand up the /scan endpoint returning the shared JSON contract. "
     "Frontend swaps from local mock to API with no UI changes."),
    ("Phase 3", "Real rPPG pipeline",
     "Replace the mock with the actual face detection · POS · HRV signal chain. "
     "Endpoint signature does not change."),
    ("Phase 4", "Docker + Cloudflare Tunnel deployment",
     "Three-container stack on the homelab, exposed via outbound tunnel. "
     "Site reachable at vitalscan.bkre8tive.com."),
    ("Phase 5", "Browser webcam capture",
     "In-browser 30-second capture using getUserMedia + MediaRecorder. "
     "Closes the loop end-to-end."),
]

top = Inches(1.55)
for phase, headline, detail in phases:
    card = filled_card(s, Inches(0.5), top, Inches(12.3), Inches(0.95))
    # Number badge
    circle = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                Inches(0.75), top + Inches(0.22),
                                Inches(0.52), Inches(0.52))
    circle.fill.solid(); circle.fill.fore_color.rgb = ACCENT
    circle.line.fill.background()
    add_text(s, Inches(0.75), top + Inches(0.24), Inches(0.52), Inches(0.5),
             phase.split()[1], size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # Headline + detail
    add_text(s, Inches(1.5), top + Inches(0.15), Inches(11.0), Inches(0.4),
             f"{phase} — {headline}", size=15, bold=True, color=NAVY)
    add_text(s, Inches(1.5), top + Inches(0.5), Inches(11.0), Inches(0.4),
             detail, size=11, color=SLATE)
    top += Inches(1.05)

set_notes(s, """BUILD PHASES (2m):

Five phases, executed in order. The critical design choice is the ordering — we're explicitly NOT building the rPPG pipeline first.

Phase 1 — React UI against mock data. This sounds backwards but it's deliberate. By building the UI first, we lock in the user experience and the JSON contract simultaneously. Everything downstream has to fit this shape.

Phase 2 — FastAPI backend returning a mock /scan response. Same contract, real HTTP service. The frontend swaps its data source from local to remote with zero UI changes. This validates the contract.

Phase 3 — now we plug in the real rPPG pipeline behind the same endpoint. If the UI worked with mock data and the contract held with HTTP, then real biomarkers slot in cleanly. No UI breakage, no contract negotiation.

Phase 4 — Dockerize and deploy through Cloudflare Tunnel. We don't deploy earlier because we want the algorithm validated locally first. Once deployed, the live URL becomes the canonical demo target.

Phase 5 — browser webcam capture. Last because it's a UI feature that depends on a working backend pipeline. Doing it last means the capture flow can be tested against real biomarker output, not mocks.

The professor's point about "design before code" maps onto this directly: we designed this ordering before writing any of the phases. The order minimizes integration risk — at each boundary, we know what works and what doesn't.""")

# ============================================================
# SLIDE 7 — DATA PLAN  [ANIMATED · 3 click reveals]
# ============================================================
s = add_slide()
header(s, "Data plan",
       "Click to step through · two datasets · one evaluation target")

groups_slide7 = []

# UBFC card
ubfc_card = filled_card(s, Inches(0.5), Inches(1.7), Inches(6.0), Inches(3.5),
                        fill=CARD_BG, border=ACCENT)
ubfc_title = add_text(s, Inches(0.75), Inches(1.85), Inches(5.5), Inches(0.5),
                     "UBFC-rPPG · real humans", size=18, bold=True, color=NAVY)
ubfc_body = add_bullets(s, Inches(0.75), Inches(2.45), Inches(5.5), Inches(2.6), [
    "42 subjects · webcam @ 30 fps",
    "Synchronized finger pulse-ox ground truth",
    "Source: Google Drive (project brief)",
    "Role: PRIMARY accuracy benchmark — this is the bar that matters",
])
groups_slide7.append([ubfc_card.shape_id, ubfc_title.shape_id, ubfc_body.shape_id])

# SCAMPS card
scamps_card = filled_card(s, Inches(6.8), Inches(1.7), Inches(6.0), Inches(3.5),
                          fill=ORANGE_BG, border=WARN)
scamps_title = add_text(s, Inches(7.05), Inches(1.85), Inches(5.5), Inches(0.5),
                       "SCAMPS · synthetic avatars", size=18, bold=True, color=NAVY)
scamps_body = add_bullets(s, Inches(7.05), Inches(2.45), Inches(5.5), Inches(2.6), [
    "10 video example subset · public",
    "Simulated PPG ground truth from avatar",
    "Source: Microsoft (danmcduff/scampsdataset)",
    "Role: SECONDARY synthetic-data baseline",
])
groups_slide7.append([scamps_card.shape_id, scamps_title.shape_id, scamps_body.shape_id])

# Evaluation target card
eval_card = filled_card(s, Inches(0.5), Inches(5.4), Inches(12.3), Inches(1.4),
                        fill=GREEN_BG, border=GOOD)
eval_label = add_text(s, Inches(0.75), Inches(5.55), Inches(11.8), Inches(0.4),
                     "EVALUATION TARGET (per rubric)", size=12, bold=True, color=GOOD)
eval_headline = add_text(s, Inches(0.75), Inches(5.95), Inches(11.8), Inches(0.5),
                        "Mean Absolute Error  <  10 BPM   on real human data",
                        size=22, bold=True, color=NAVY)
eval_caption = add_text(s, Inches(0.75), Inches(6.45), Inches(11.8), Inches(0.3),
                       "We compare against the CHROM baseline (de Haan & Jeanne 2013) on the same input — like-for-like comparison.",
                       size=11, color=SLATE)
groups_slide7.append([eval_card.shape_id, eval_label.shape_id,
                      eval_headline.shape_id, eval_caption.shape_id])

add_sequential_fade(s, groups_slide7)

set_notes(s, """DATA PLAN (90s):

Two datasets — real and synthetic — and one evaluation harness that reads both.

UBFC-rPPG is our primary dataset because it's real humans with real ground truth from a synchronized finger pulse-oximeter. Forty-two subjects. If we're going to claim our pipeline works on people, UBFC is where we prove it.

SCAMPS is the secondary dataset — Microsoft's synthetic facial avatars. Useful as a synthetic-data baseline and stress test, but synthetic skin reflectance doesn't perfectly replicate real capillary blood-volume changes, so we don't treat it as the real bar.

The bottom green callout is the rubric target — MAE under 10 BPM on real human data. That's our pass/fail criterion. We'll also report against CHROM, which is the canonical pre-POS classical algorithm from 2013, so the comparison is like-for-like.

The professor specifically said "the system is as good as the data" — that's why we're being explicit about which dataset is primary and which is secondary. We're not going to game the average by mixing them.""")

# ============================================================
# SLIDE 8 — RISKS & MITIGATIONS
# ============================================================
s = add_slide()
header(s, "Risks & mitigations",
       "What could go wrong, and what we'll do about it")

risks = [
    ("Blood pressure can't be derived from classical rPPG",
     "Physical limitation — needs cuff calibration per subject",
     "Return mock BP flagged as known limitation; document in writeup",
     DANGER, RED_BG),
    ("Head motion + bad lighting destroy the signal",
     "RGB variation from motion dwarfs the 1% pulse signal",
     "Forehead + cheek ROIs (less motion); bandpass filter; POS projection",
     WARN, ORANGE_BG),
    ("Synthetic SCAMPS data has facial-animation artifacts",
     "Avatar animation frequencies sit on top of the pulse band",
     "Tighten FFT search to 60–210 BPM; UBFC remains primary benchmark",
     WARN, ORANGE_BG),
    ("Deep-learning extension is a stretch goal",
     "Classical pipeline must clear rubric first; DL is opt-in",
     "Scope DL out of Phase 1–5; scaffolding only; revisit post-rubric",
     ACCENT, CARD_BG),
]

top = Inches(1.55)
for risk, why, mit, border_color, bg_color in risks:
    card = filled_card(s, Inches(0.5), top, Inches(12.3), Inches(1.20),
                       fill=bg_color, border=border_color)
    add_text(s, Inches(0.75), top + Inches(0.13), Inches(11.8), Inches(0.4),
             f"RISK — {risk}", size=13, bold=True, color=border_color)
    add_text(s, Inches(0.75), top + Inches(0.52), Inches(11.8), Inches(0.3),
             f"Why: {why}", size=11, color=SLATE)
    add_text(s, Inches(0.75), top + Inches(0.85), Inches(11.8), Inches(0.3),
             f"Mitigation: {mit}", size=11, bold=True, color=NAVY)
    top += Inches(1.30)

set_notes(s, """RISKS & MITIGATIONS (90s):

Four risks we've identified up-front. Each one named, each one mitigated.

First — the big one — blood pressure. Classical rPPG can't derive BP without a cuff-calibrated reference. This is a physical limit, not an engineering problem. We're going to return a flagged placeholder value and document it as a known limitation in the writeup. The honest answer is "this is what classical rPPG can do; getting BP would require a different technique."

Second — head motion and lighting. The pulse signal is one percent of reflected light; head motion is much bigger. Three mitigations stack: we pick forehead and cheeks because they move less, we bandpass-filter the signal to the heart-rate band, and the POS algorithm itself is designed to suppress motion-correlated noise.

Third — SCAMPS synthetic-data artifacts. The synthetic avatars have facial animation that puts strong artifacts in the low-frequency band where slow heart rates live. We're handling this by tightening the FFT search window to 60–210 BPM and treating UBFC as the primary benchmark — synthetic data is for stress-testing, not for the headline number.

Fourth — the deep-learning extension was a stretch deliverable. We're explicitly scoping it out of the main five phases. If we have time after the classical pipeline clears the rubric, we'll come back to it. If we don't, we don't — the rubric doesn't require it.

The pattern: every risk has a named mitigation. That's what "designing before coding" means in practice.""")

# ============================================================
# SLIDE 9 — TIMELINE / NEXT STEPS
# ============================================================
s = add_slide()
header(s, "Timeline & next steps",
       "What we're doing this week, what comes after")

# Two columns: now vs week 8
filled_card(s, Inches(0.5), Inches(1.7), Inches(6.0), Inches(5.0),
            fill=CARD_BG, border=ACCENT)
add_text(s, Inches(0.75), Inches(1.85), Inches(5.5), Inches(0.4),
         "NOW — design checkpoint", size=12, bold=True, color=ACCENT)
add_text(s, Inches(0.75), Inches(2.3), Inches(5.5), Inches(0.5),
         "What's done in planning", size=18, bold=True, color=NAVY)
add_bullets(s, Inches(0.75), Inches(2.9), Inches(5.5), Inches(3.6), [
    "API contract — defined and agreed",
    "System architecture — three layers locked",
    "Algorithm — POS algorithm chosen with rationale",
    "Tech stack — every choice has a reason",
    "Build phases — ordered to minimize integration risk",
    "Risks — named with mitigations",
])

filled_card(s, Inches(6.8), Inches(1.7), Inches(6.0), Inches(5.0),
            fill=GREEN_BG, border=GOOD)
add_text(s, Inches(7.05), Inches(1.85), Inches(5.5), Inches(0.4),
         "NEXT — execution toward Week 8", size=12, bold=True, color=GOOD)
add_text(s, Inches(7.05), Inches(2.3), Inches(5.5), Inches(0.5),
         "What we execute next", size=18, bold=True, color=NAVY)
add_bullets(s, Inches(7.05), Inches(2.9), Inches(5.5), Inches(3.6), [
    "Phase 1 — React UI against mock JSON",
    "Phase 2 — FastAPI mock /scan endpoint",
    "Phase 3 — real rPPG pipeline (face → POS → HRV)",
    "Phase 4 — Docker stack + Cloudflare Tunnel",
    "Phase 5 — browser webcam capture",
    "Week 8 — group results presentation",
])

set_notes(s, """TIMELINE / NEXT STEPS (60s):

Left column — what we've finished in the design phase. Six items. We've defined the API contract that other groups can consume, locked in the three-layer architecture, chosen the POS algorithm over alternatives (with rationale), justified every tech-stack pick, ordered our build phases to minimize integration risk, and named our risks with concrete mitigations.

That's the planning deliverable. None of it is code. All of it is the foundation that makes the code straightforward.

Right column — execution. Five phases in order. By Week 8 we'll come back with the results — accuracy numbers on UBFC, the live deployment, the integration with the other groups. That's a different presentation.

The professor's point about being systems engineers, not coders — the left column is what systems engineers produce. The right column is what coders execute after the design is solid.""")

# ============================================================
# SLIDE 10 — CLOSE
# ============================================================
s = add_slide()
panel = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
panel.line.fill.background()
panel.fill.solid()
panel.fill.fore_color.rgb = NAVY

add_text(s, Inches(0.6), Inches(0.8), Inches(12), Inches(0.5),
         "AIT 500 · GROUP 1 · PLANNING CHECKPOINT", size=12, bold=True, color=ACCENT)
add_text(s, Inches(0.6), Inches(1.6), Inches(12), Inches(1.3),
         "Designed before coding.", size=54, bold=True, color=WHITE)
add_text(s, Inches(0.6), Inches(3.0), Inches(12), Inches(0.6),
         "Questions about the blueprint?", size=24, color=WHITE)

# anchor points
add_text(s, Inches(0.6), Inches(4.2), Inches(12), Inches(0.4),
         "DESIGN SUMMARY", size=14, bold=True, color=ACCENT)
add_bullets(s, Inches(0.6), Inches(4.7), Inches(12), Inches(2.2), [
    "Three-layer architecture: client / Cloudflare / homelab",
    "Four-stage rPPG pipeline: face → ROI → POS → BPM",
    "Five build phases, ordered to validate each boundary",
    "Target: MAE < 10 BPM on UBFC-rPPG (real humans)",
], size=15, color=WHITE)

add_text(s, Inches(0.6), Inches(6.8), Inches(12), Inches(0.3),
         "Group 1 — Germaine · Jason · Daray · Abner",
         size=12, color=ACCENT)

set_notes(s, """CLOSE (30s):

To summarize in one line: we designed before we coded.

The architecture is three layers — client, Cloudflare edge, homelab backend. The algorithm is the four-stage POS-based signal chain. We have five build phases in a specific order that validates each boundary as we go. And we're aiming for under 10 BPM MAE on the real-human UBFC dataset, which is the rubric target.

Happy to take questions on any of the design choices. Thank you.""")


# write
prs.save(OUT)
print(f"wrote {OUT}")
print(f"  slides: {len(prs.slides)}")
