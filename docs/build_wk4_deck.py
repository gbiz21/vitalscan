"""Build the Week 4 submission deck: Model Development and Framework Design.

Course:      AIT 500 · Westcliff University · MS in AI
Assignment:  Wk 4 — Project: Model Development and Framework Design (CLO 2/3/4)
Output:      docs/Wk4_VitalScan_Model_Development_and_Framework_Design.pptx

Per the Moodle brief:
  "Begin designing your AI model or framework. This involves outlining the
   architecture, defining algorithms, and determining the data requirements
   for your model.
   Deliverable: A design document that includes the model architecture,
   chosen algorithms, and data flow diagrams illustrating how your model
   will function."
"""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parent.parent
FIG_DIR = ROOT / "docs" / "figures"
OUT = ROOT / "docs" / "Wk4_VitalScan_Model_Development_and_Framework_Design.pptx"

# --- palette (consistent with prior decks) ---
NAVY = RGBColor(0x0F, 0x1F, 0x44)
SLATE = RGBColor(0x33, 0x44, 0x5C)
INK = RGBColor(0x1A, 0x1A, 0x1A)
MUTED = RGBColor(0x55, 0x66, 0x7A)
ACCENT = RGBColor(0x1F, 0x6F, 0xEB)
GOOD = RGBColor(0x1E, 0x88, 0x4F)
WARN = RGBColor(0xC0, 0x73, 0x1D)
DANGER = RGBColor(0xC0, 0x39, 0x2B)
SOFT_BG = RGBColor(0xF6, 0xF8, 0xFB)
CARD_BG = RGBColor(0xEE, 0xF3, 0xFA)
GREEN_BG = RGBColor(0xE9, 0xF5, 0xEC)
ORANGE_BG = RGBColor(0xFD, 0xF1, 0xE3)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


# ---------- helpers ----------
def add_slide():
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE
    return s


def add_text(slide, left, top, width, height, text, *, size=18, bold=False,
             color=INK, align=PP_ALIGN.LEFT, font="Calibri"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.name = font; r.font.size = Pt(size)
        r.font.bold = bold; r.font.color.rgb = color
    return tb


def add_bullets(slide, left, top, width, height, bullets, *, size=12, color=INK):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(4)
        r = p.add_run()
        r.text = f"•  {b}"
        r.font.name = "Calibri"
        r.font.size = Pt(size)
        r.font.color.rgb = color
    return tb


def filled_card(slide, left, top, width, height, *,
                fill=SOFT_BG, border=ACCENT, radius=0.06):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  left, top, width, height)
    card.adjustments[0] = radius
    card.fill.solid()
    card.fill.fore_color.rgb = fill
    card.line.color.rgb = border
    card.line.width = Pt(1.25)
    card.shadow.inherit = False
    return card


def header(slide, title, subtitle=None):
    stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, Inches(0.18))
    stripe.line.fill.background()
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = ACCENT
    add_text(slide, Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.6),
             title, size=26, bold=True, color=NAVY)
    if subtitle:
        add_text(slide, Inches(0.5), Inches(0.88), Inches(12.3), Inches(0.4),
                 subtitle, size=13, color=MUTED)
    add_text(slide, Inches(0.5), Inches(7.05), Inches(12.3), Inches(0.3),
             "AIT 500 · Westcliff University · Group 1 · Week 4 Submission",
             size=10, color=MUTED)


def set_notes(slide, text):
    try:
        notes_tf = slide.notes_slide.notes_text_frame
        if notes_tf is not None:
            notes_tf.text = text
    except (AttributeError, ValueError):
        pass


def add_table(slide, left, top, col_widths, rows, *, header_row=True,
              header_fill=NAVY, header_fg=WHITE, body_fg=INK,
              row_h=Inches(0.45), header_h=Inches(0.45), font_size=11):
    """Render a clean styled table with header + body rows."""
    y = top
    for r_idx, row in enumerate(rows):
        x = left
        is_header = header_row and r_idx == 0
        h = header_h if is_header else row_h
        for c_idx, val in enumerate(row):
            w = col_widths[c_idx]
            cell = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
            if is_header:
                cell.fill.solid(); cell.fill.fore_color.rgb = header_fill
                cell.line.color.rgb = WHITE
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = SOFT_BG if r_idx % 2 == 0 else WHITE
                cell.line.color.rgb = RGBColor(0xD0, 0xD7, 0xDF)
            cell.line.width = Pt(0.5)
            add_text(slide, x + Inches(0.1), y + Inches(0.08),
                     w - Inches(0.2), h - Inches(0.15),
                     val, size=font_size,
                     bold=is_header,
                     color=header_fg if is_header else body_fg)
            x += w
        y += h
    return y


# ============================================================
# SLIDE 1 — TITLE
# ============================================================
s = add_slide()
panel = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(5.6), SH)
panel.line.fill.background()
panel.fill.solid()
panel.fill.fore_color.rgb = NAVY

add_text(s, Inches(0.6), Inches(0.55), Inches(4.8), Inches(0.4),
         "AIT 500 · WESTCLIFF UNIVERSITY", size=11, bold=True, color=ACCENT)
add_text(s, Inches(0.6), Inches(0.95), Inches(4.8), Inches(0.4),
         "M.S. ARTIFICIAL INTELLIGENCE", size=10, color=WHITE)
add_text(s, Inches(0.6), Inches(1.6), Inches(5.0), Inches(1.4),
         "VitalScan", size=60, bold=True, color=WHITE)
add_text(s, Inches(0.6), Inches(3.05), Inches(5.0), Inches(1.0),
         "Model Development &\nFramework Design",
         size=22, color=WHITE)
add_text(s, Inches(0.6), Inches(4.7), Inches(4.8), Inches(0.4),
         "GROUP 1 · rPPG SIGNAL EXTRACTION",
         size=11, bold=True, color=ACCENT)
add_text(s, Inches(0.6), Inches(5.4), Inches(4.8), Inches(0.4),
         "WEEK 4 · CLO 2 · CLO 3 · CLO 4",
         size=10, bold=True, color=ACCENT)
add_text(s, Inches(0.6), Inches(6.6), Inches(4.8), Inches(0.4),
         "vitalscan.bkre8tive.com", size=12, bold=True, color=ACCENT)

# right side — submission card
add_text(s, Inches(6.2), Inches(0.85), Inches(6.5), Inches(0.4),
         "WEEK 4 SUBMISSION", size=11, bold=True, color=MUTED)
add_text(s, Inches(6.2), Inches(1.3), Inches(6.5), Inches(0.7),
         "Project: Model Development\nand Framework Design",
         size=26, bold=True, color=NAVY)

def assignment_card(top, label, body):
    card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                              Inches(6.2), top, Inches(6.3), Inches(0.85))
    card.adjustments[0] = 0.18
    card.fill.solid(); card.fill.fore_color.rgb = CARD_BG
    card.line.color.rgb = ACCENT; card.line.width = Pt(1)
    add_text(s, Inches(6.45), top + Inches(0.12),
             Inches(2.6), Inches(0.4), label, size=13, bold=True, color=NAVY)
    add_text(s, Inches(6.45), top + Inches(0.45),
             Inches(5.7), Inches(0.4), body, size=11, color=SLATE)

assignment_card(Inches(2.65), "Deliverable 1", "Model architecture")
assignment_card(Inches(3.60), "Deliverable 2", "Chosen algorithms with rationale")
assignment_card(Inches(4.55), "Deliverable 3", "Data flow diagrams")
assignment_card(Inches(5.50), "Reference",      "Bawack et al. (2021) AI research framework")

add_text(s, Inches(6.2), Inches(6.55), Inches(6.5), Inches(0.4),
         "Germaine Beazer · Jason · Daray · Abner",
         size=11, color=MUTED)

set_notes(s, """Title slide for the Week 4 submission. Three deliverables on the right correspond directly to the assignment brief: (1) model architecture, (2) chosen algorithms with rationale, (3) data flow diagrams. Supplementary reading is Bawack et al. 2021. Submission for AIT 500 / CLO 2 / CLO 3 / CLO 4. Live deployment at vitalscan.bkre8tive.com.""")

# ============================================================
# SLIDE 2 — ASSIGNMENT CONTEXT + FRAMEWORK REFERENCE
# ============================================================
s = add_slide()
header(s, "Assignment context & framework reference",
       "Wk 4 — Project: Model Development and Framework Design · CLO 2 / 3 / 4")

# Left — direct brief quote
card_l = filled_card(s, Inches(0.5), Inches(1.5), Inches(6.0), Inches(5.0))
add_text(s, Inches(0.75), Inches(1.7), Inches(5.5), Inches(0.4),
         "ASSIGNMENT — verbatim from Moodle", size=11, bold=True, color=ACCENT)
add_text(s, Inches(0.75), Inches(2.15), Inches(5.5), Inches(3.6),
         "\"Begin designing your AI model or framework.\n"
         "This involves outlining the architecture, defining\n"
         "algorithms, and determining the data requirements\n"
         "for your model.\n\n"
         "Deliverable: A design document that includes the\n"
         "model architecture, chosen algorithms, and data\n"
         "flow diagrams illustrating how your model will function.\"",
         size=12, color=INK)

add_text(s, Inches(0.75), Inches(5.85), Inches(5.5), Inches(0.4),
         "CLOs ADDRESSED", size=11, bold=True, color=MUTED)
add_bullets(s, Inches(0.75), Inches(6.2), Inches(5.5), Inches(0.85), [
    "CLO 2 — AI fundamentals: types, methods, and limits",
    "CLO 3 — Algorithm selection with rationale",
    "CLO 4 — Applied AI: design + implementation",
], size=11)

# Right — Bawack et al. framework
card_r = filled_card(s, Inches(6.8), Inches(1.5), Inches(6.0), Inches(5.0),
                     fill=CARD_BG, border=ACCENT)
add_text(s, Inches(7.05), Inches(1.7), Inches(5.5), Inches(0.4),
         "FRAMEWORK REFERENCE", size=11, bold=True, color=ACCENT)
add_text(s, Inches(7.05), Inches(2.1), Inches(5.5), Inches(0.4),
         "Bawack et al. (2021)", size=14, bold=True, color=NAVY)
add_text(s, Inches(7.05), Inches(2.5), Inches(5.5), Inches(0.5),
         "\"A framework for understanding artificial intelligence\n"
         "research: insights from practice\"",
         size=11, color=SLATE)
add_text(s, Inches(7.05), Inches(3.2), Inches(5.5), Inches(0.3),
         "J. Enterprise Information Management, 3(2), 645–678",
         size=10, color=MUTED)

add_text(s, Inches(7.05), Inches(3.75), Inches(5.5), Inches(0.4),
         "VITALSCAN, MAPPED TO THE FRAMEWORK", size=11, bold=True, color=ACCENT)
add_bullets(s, Inches(7.05), Inches(4.15), Inches(5.5), Inches(2.4), [
    "AI capability — computer vision + signal extraction",
    "Method family — hybrid: pretrained CNN + classical SP",
    "Application domain — health / contactless biomarkers",
    "Data type — short facial video (RGB image sequence)",
    "Validation — supervised held-out on labeled data (UBFC)",
    "Deployment — production REST API, live containerized stack",
], size=11)

set_notes(s, """The assignment brief is reproduced verbatim on the left. The right column ties our work to the Bawack et al. (2021) framework — they propose a typology of AI research along several dimensions, and VitalScan maps cleanly to each: we are doing computer-vision-based signal extraction (capability), using a hybrid CNN + classical signal-processing method (method family), applied to health biomarkers (domain), with video as the data type, supervised held-out evaluation, and a live production deployment.""")

# ============================================================
# SLIDE 3 — PROBLEM & OBJECTIVE
# ============================================================
s = add_slide()
header(s, "Problem statement & objective",
       "What the AI model is designed to do, and why")

filled_card(s, Inches(0.5), Inches(1.55), Inches(6.0), Inches(2.6))
add_text(s, Inches(0.75), Inches(1.7), Inches(5.5), Inches(0.4),
         "INPUT", size=11, bold=True, color=MUTED)
add_text(s, Inches(0.75), Inches(2.1), Inches(5.5), Inches(0.5),
         "30-second facial video", size=20, bold=True, color=NAVY)
add_bullets(s, Inches(0.75), Inches(2.7), Inches(5.5), Inches(1.4), [
    "Standard webcam or phone camera · 30 fps",
    "No special hardware · no contact sensor",
    "Cooperative subject, mostly frontal, well-lit",
], size=12)

filled_card(s, Inches(6.8), Inches(1.55), Inches(6.0), Inches(2.6),
            fill=GREEN_BG, border=GOOD)
add_text(s, Inches(7.05), Inches(1.7), Inches(5.5), Inches(0.4),
         "OUTPUT", size=11, bold=True, color=GOOD)
add_text(s, Inches(7.05), Inches(2.1), Inches(5.5), Inches(0.5),
         "Biomarker JSON contract", size=20, bold=True, color=NAVY)
add_bullets(s, Inches(7.05), Inches(2.7), Inches(5.5), Inches(1.4), [
    "Heart rate (BPM) · HRV SDNN (ms) · stress (0–1)",
    "Blood pressure placeholder (flagged with low confidence)",
    "Each biomarker carries {value, confidence, unit}",
], size=12)

add_text(s, Inches(0.5), Inches(4.4), Inches(12.3), Inches(0.4),
         "WHY THIS PROBLEM IS AN AI PROBLEM (CLO 2)",
         size=12, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

filled_card(s, Inches(0.5), Inches(4.85), Inches(12.3), Inches(1.85))
add_bullets(s, Inches(0.75), Inches(5.0), Inches(11.8), Inches(1.6), [
    "Perception under noise — pulse signal is 0.1–1 % of reflected light, dominated by lighting and motion artifacts",
    "Pattern recognition — a learned CNN identifies and tracks the face; classical projection isolates the periodic pulse signature",
    "Generalization required — the model must work on unseen subjects across skin tone, lighting, and camera variation",
    "Quantitative output — model produces a continuous heart-rate value with a confidence estimate, not a categorical label",
], size=12)

set_notes(s, """Problem and objective statement. The input is a 30-second facial video; the output is a biomarker JSON object with four physiological biomarkers, each annotated with a confidence score the downstream consumer can use to weight risk. The lower half explains why this qualifies as an AI problem under CLO 2: perception under heavy noise, pattern recognition through learned and classical methods, generalization to unseen subjects, and quantitative continuous output rather than categorical classification.""")

# ============================================================
# SLIDE 4 — MODEL ARCHITECTURE
# ============================================================
s = add_slide()
header(s, "The AI model — architecture",
       "Deliverable 1: hybrid model with pretrained CNN frontend + classical signal-processing backend")

# Two column layout: left text, right schematic
add_text(s, Inches(0.5), Inches(1.55), Inches(6.0), Inches(0.5),
         "Architectural choice — hybrid", size=14, bold=True, color=NAVY)
add_bullets(s, Inches(0.5), Inches(2.05), Inches(6.0), Inches(4.7), [
    "Stage 1 — Pretrained deep CNN (MediaPipe FaceMesh, Google)\n"
    "      Trained on millions of faces; emits 468 landmarks / frame",
    "Stage 2 — ROI extraction (polygon mask over forehead + cheeks)\n"
    "      Compute mean RGB inside each region per frame",
    "Stage 3 — POS algorithm (Wang et al. 2017)\n"
    "      Linear projection isolates pulse from lighting / motion",
    "Stage 4 — Butterworth bandpass (0.7–4 Hz, order 3)\n"
    "      Zero-phase filter keeps only the heart-rate band",
    "Stage 5 — FFT + peak in 60–210 BPM band\n"
    "      Dominant frequency → heart rate (BPM)",
    "Stage 6 — Peak detection + IBI → SDNN + LF/HF\n"
    "      Same waveform yields HRV and stress index",
    "Stage 7 — Confidence scoring per biomarker\n"
    "      Output: {value, confidence, unit} for each biomarker",
], size=11)

# Right side — actual architecture figure
s.shapes.add_picture(str(FIG_DIR / "figure1-signal-chain.png"),
                     Inches(6.8), Inches(1.55), width=Inches(6.0))

add_text(s, Inches(6.8), Inches(6.6), Inches(6.0), Inches(0.3),
         "Figure 1 — Signal chain (rPPG pipeline)",
         size=10, color=MUTED, align=PP_ALIGN.CENTER)

set_notes(s, """Deliverable 1 — model architecture. The model is a HYBRID of one pretrained deep CNN (MediaPipe FaceMesh) followed by six classical signal-processing stages. The CNN stage demonstrates we are using a learned model — it is exactly the kind of multi-layer feature-learning network shown in the cat-classification example from the 2026-05-27 class lecture. The classical stages are explicit, mathematically defined transforms with documented inputs and outputs (defended in our Week 3 research deliverable). Confidence scoring (added 2026-05-25 per professor feedback) ensures every output carries a reliability estimate.""")

# ============================================================
# SLIDE 5 — CHOSEN ALGORITHMS (Deliverable 2)
# ============================================================
s = add_slide()
header(s, "Chosen algorithms with rationale",
       "Deliverable 2: each stage names the algorithm, its alternatives, and the selection reason")

col_w = [Inches(2.5), Inches(2.7), Inches(3.5), Inches(3.5)]
rows = [
    ["Stage", "Chosen algorithm", "Alternatives considered", "Rationale"],
    ["Face detection",  "MediaPipe FaceMesh",
     "Haar, dlib HOG, MTCNN, RetinaFace",
     "468 dense landmarks; CPU-realtime; pretrained"],
    ["ROI extraction",  "Forehead + cheek polygons",
     "Bounding box, full face, lower face",
     "Highest capillary density; lowest motion"],
    ["Pulse extraction", "POS (Wang 2017)",
     "GREEN, ICA, CHROM, PBV, deep-learning",
     "Strongest classical; interpretable waveform"],
    ["Bandpass",        "Butterworth, order 3, filtfilt",
     "Moving average, FIR, Chebyshev",
     "Zero-phase, sharp roll-off, scipy-tested"],
    ["HR estimation",   "FFT + peak in 60–210 BPM",
     "Wavelets, autocorrelation, time-domain",
     "Robust on 30-s clips; standard in field"],
    ["HRV",             "SDNN (time domain)",
     "RMSSD, pNN50, triangular index",
     "Most widely reported clinical metric"],
    ["Stress index",    "LF/HF ratio → sigmoid",
     "Baevsky index, SDNN threshold",
     "Frequency-domain ANS standard"],
    ["Confidence",      "Peak prominence + beat count + IBI CV",
     "Bayesian uncertainty, ensemble variance",
     "Math-based, no second model required"],
]
add_table(s, Inches(0.5), Inches(1.55), col_w, rows,
          row_h=Inches(0.55), header_h=Inches(0.45), font_size=10)

set_notes(s, """Deliverable 2 — chosen algorithms with rationale. Each row is one pipeline stage: the algorithm selected, the alternatives we considered, and why this one. This table satisfies the assignment's "defining algorithms" requirement. Full reasoning for each row is documented in the Week 3 research deliverable (docs/Week3_Research_AI_Method_Selection.md), which surveys five classical rPPG algorithms and five deep-learning rPPG architectures before justifying POS + MediaPipe as the chosen primary path.""")

# ============================================================
# SLIDE 6 — DATA FLOW DIAGRAM (Deliverable 3)
# ============================================================
s = add_slide()
header(s, "Data flow diagram — model end-to-end",
       "Deliverable 3: tensor-level data flow with shapes annotated at each stage")

# Flow diagram built from boxes + arrows. Vertical flow.
# Each step: a card with shape on left, transform name on right.
FLOW_STAGES = [
    ("Video",          "(900 frames, 1080, 1920, 3)", "uint8 RGB · 30 s @ 30 fps", ACCENT, CARD_BG),
    ("MediaPipe CNN",  "(900, 468, 3)",                "pretrained face-mesh landmarks", ACCENT, CARD_BG),
    ("ROI polygon mask + mean RGB", "(900, 3)",        "mean RGB averaged over forehead + cheeks", ACCENT, CARD_BG),
    ("POS algorithm (1.6 s window)", "(900,)",         "linear projection isolates pulse signal", GOOD, GREEN_BG),
    ("Butterworth bandpass 0.7–4 Hz", "(900,)",        "scipy.signal.filtfilt · zero-phase", GOOD, GREEN_BG),
    ("FFT + peak (60–210 BPM)",      "scalar BPM",     "+ confidence from peak-to-mean ratio", WARN, ORANGE_BG),
    ("Peak detect → IBIs → SDNN, Welch LF/HF", "scalar ms · scalar [0,1]", "+ confidence from beat-count + IBI plausibility", WARN, ORANGE_BG),
    ("Biomarker JSON contract",      "{value, confidence, unit} ×4", "shared output to Groups 3 + 4", DANGER, RGBColor(0xFA, 0xEC, 0xEA)),
]

x_card = Inches(0.5); w_card = Inches(12.3); h_card = Inches(0.55)
y = Inches(1.55)
gap = Inches(0.10)
for stage, shape_text, detail, border, fill in FLOW_STAGES:
    card = filled_card(s, x_card, y, w_card, h_card, fill=fill, border=border)
    add_text(s, x_card + Inches(0.3), y + Inches(0.13), Inches(4.5), Inches(0.4),
             stage, size=13, bold=True, color=NAVY)
    add_text(s, x_card + Inches(5.0), y + Inches(0.13), Inches(3.0), Inches(0.4),
             shape_text, size=11, bold=True, color=border,
             font="Consolas")
    add_text(s, x_card + Inches(8.2), y + Inches(0.16), Inches(4.0), Inches(0.4),
             detail, size=10, color=SLATE)
    y += h_card + gap

# Tiny legend
add_text(s, Inches(0.5), Inches(6.85), Inches(12.3), Inches(0.3),
         "Blue = ingest/preprocessing  ·  Green = pulse extraction  ·  Amber = biomarker derivation  ·  Red = packaged output",
         size=10, color=MUTED, align=PP_ALIGN.CENTER)

set_notes(s, """Deliverable 3 — data flow diagram. Vertical flow shows the exact tensor shape at each stage, what transform produces it, and the implementation detail. Color bands group the stages: ingest + preprocessing in blue, pulse extraction in green, biomarker derivation in amber, packaged output in red. This is the diagram the assignment explicitly asks for — "data flow diagrams illustrating how your model will function."

Key numbers a reviewer can verify: 30 fps × 30 s = 900 frames; MediaPipe outputs 468 landmarks each with (x, y, z); after ROI masking + averaging we collapse to a single (R, G, B) per frame; the entire heavy lifting reduces a 5+ GB raw video down to four scalar biomarkers.""")

# ============================================================
# SLIDE 7 — DATA REQUIREMENTS
# ============================================================
s = add_slide()
header(s, "Data requirements",
       "What the model needs at inference time, and what the ground-truth data must contain to evaluate it")

# Two side-by-side tables: inference input vs evaluation data
add_text(s, Inches(0.5), Inches(1.55), Inches(6.0), Inches(0.4),
         "INFERENCE INPUT (per scan)", size=12, bold=True, color=ACCENT)
rows1 = [
    ["Requirement", "Specification"],
    ["Format",      ".mp4 / .webm / .mov / .avi / .mkv"],
    ["Codec",       "H.264, VP8/9 (extension-permissive)"],
    ["Duration",    "30 s preferred (min 5 s)"],
    ["Frame rate",  "30 fps (works 24–60)"],
    ["Resolution",  "≥360p; face must be detectable"],
    ["Color",       "RGB (3 channels)"],
    ["Subject",     "1 cooperative face, mostly frontal, lit"],
    ["Max upload",  "100 MB"],
]
add_table(s, Inches(0.5), Inches(2.0), [Inches(2.0), Inches(4.0)],
          rows1, row_h=Inches(0.4), header_h=Inches(0.4), font_size=11)

add_text(s, Inches(6.8), Inches(1.55), Inches(6.0), Inches(0.4),
         "GROUND-TRUTH DATA (per evaluation subject)",
         size=12, bold=True, color=GOOD)
rows2 = [
    ["Field", "Format / range"],
    ["Video",          "Same as inference input"],
    ["PPG waveform",   "Sampled trace from finger pulse-oximeter"],
    ["Per-frame HR",   "BPM, sample-synchronized with video"],
    ["Timestamps",     "Seconds, scientific notation"],
    ["Subject ID",     "Integer (used for train/test split)"],
    ["Lighting note",  "Recorded environment metadata"],
    ["UBFC schema",    "video.avi + ground_truth.txt (3 rows)"],
    ["SCAMPS schema",  "Simulated PPG (synthetic baseline)"],
]
add_table(s, Inches(6.8), Inches(2.0), [Inches(2.0), Inches(4.0)],
          rows2, row_h=Inches(0.4), header_h=Inches(0.4), font_size=11)

filled_card(s, Inches(0.5), Inches(5.85), Inches(12.3), Inches(0.95),
            fill=CARD_BG, border=ACCENT)
add_text(s, Inches(0.75), Inches(5.95), Inches(11.8), Inches(0.4),
         "DATASETS — already secured", size=11, bold=True, color=ACCENT)
add_bullets(s, Inches(0.75), Inches(6.3), Inches(11.8), Inches(0.5), [
    "UBFC-rPPG — 42 real human subjects · primary benchmark (Univ. de Bourgogne, via Google Drive)",
    "SCAMPS — 10 synthetic avatar videos · secondary baseline (Microsoft Research, public subset)",
], size=11)

set_notes(s, """Data requirements split into two halves. Left — what the model needs to run inference on a new video: file format, codec, duration, frame rate, resolution, color, subject conditions, and the 100 MB upload cap. Right — what the ground-truth dataset must contain so we can evaluate model accuracy: the paired video, the PPG waveform, per-frame heart-rate labels, timestamps, and subject identifiers for the train/test split. Both UBFC-rPPG and SCAMPS are already secured and integrated into the evaluation harness.""")

# ============================================================
# SLIDE 8 — CONFIDENCE SCORING
# ============================================================
s = add_slide()
header(s, "Confidence scoring per biomarker",
       "Model-development update — every output carries a {value, confidence, unit} object so downstream consumers can weight by reliability")

# Three cards: HR, HRV/stress, BP
def conf_card(left, top, color, bg, name, formula, description):
    filled_card(s, left, top, Inches(4.0), Inches(3.6), fill=bg, border=color)
    add_text(s, left + Inches(0.2), top + Inches(0.15), Inches(3.6), Inches(0.4),
             name, size=14, bold=True, color=color)
    add_text(s, left + Inches(0.2), top + Inches(0.6), Inches(3.6), Inches(0.6),
             "Confidence formula", size=10, bold=True, color=MUTED)
    add_text(s, left + Inches(0.2), top + Inches(0.95), Inches(3.6), Inches(0.7),
             formula, size=10, color=INK, font="Consolas")
    add_text(s, left + Inches(0.2), top + Inches(1.95), Inches(3.6), Inches(0.4),
             "What it captures", size=10, bold=True, color=MUTED)
    add_text(s, left + Inches(0.2), top + Inches(2.30), Inches(3.6), Inches(1.2),
             description, size=10, color=INK)

conf_card(Inches(0.5), Inches(1.7), ACCENT, CARD_BG,
          "Heart rate",
          "conf = clip((peak/mean - 1.5) / 4, 0, 1)",
          "Ratio of FFT-band peak magnitude to mean magnitude. Sharp dominant peak → high confidence; flat noisy spectrum → low.")

conf_card(Inches(4.7), Inches(1.7), GOOD, GREEN_BG,
          "HRV (SDNN) + Stress",
          "conf = min(1, n_peaks/N) × ibi_penalty",
          "Beat count vs target N (20 for HRV, 25 for stress). IBI coefficient of variation > 0.30 triggers a penalty for likely missed/extra beats.")

conf_card(Inches(8.9), Inches(1.7), WARN, ORANGE_BG,
          "Blood pressure",
          "conf = 0.10 (fixed)",
          "Low floor by design. Classical rPPG physically cannot derive BP without per-subject cuff calibration; placeholder value carries a note flag.")

# Bottom note
add_text(s, Inches(0.5), Inches(5.55), Inches(12.3), Inches(0.4),
         "WHY THIS MATTERS FOR THE MODEL CONTRACT",
         size=11, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
filled_card(s, Inches(0.5), Inches(5.95), Inches(12.3), Inches(1.0))
add_bullets(s, Inches(0.75), Inches(6.05), Inches(11.8), Inches(0.9), [
    "Group 3 (risk classification) weights each biomarker contribution by its confidence — high-confidence reading drives risk more than low",
    "Low confidence + note field communicates known limits (BP placeholder) without breaking the JSON contract or guessing a value",
    "Confidence emerges from the math itself (signal SNR, beat count, IBI plausibility) — no second model to train, no additional latency",
], size=11)

set_notes(s, """Confidence scoring is the model-development update added 2026-05-25 per professor feedback. Each scalar biomarker has a confidence formula derived from signal-quality metrics already computed by the pipeline — no extra model required. Heart-rate confidence is the FFT peak prominence ratio; HRV/stress confidence is beat count modulated by IBI plausibility; BP carries a fixed low confidence with an explanatory note because classical rPPG cannot derive BP. The contract change lets Group 3 weight risk assessments correctly: a high-confidence elevated heart rate should drive the recommendation more strongly than a low-confidence one.""")

# ============================================================
# SLIDE 9 — VALIDATION PLAN + EARLY RESULTS
# ============================================================
s = add_slide()
header(s, "Validation plan & early results",
       "Held-out evaluation against pulse-oximeter ground truth · classical method has no training step, so the whole UBFC set is held-out")

# Validation plan card
filled_card(s, Inches(0.5), Inches(1.55), Inches(6.0), Inches(2.5))
add_text(s, Inches(0.75), Inches(1.7), Inches(5.5), Inches(0.4),
         "VALIDATION STRATEGY", size=11, bold=True, color=ACCENT)
add_bullets(s, Inches(0.75), Inches(2.1), Inches(5.5), Inches(2.0), [
    "Primary benchmark — full UBFC-rPPG (n = 42 real subjects)",
    "Metric 1 — MAE in BPM (primary, rubric target)",
    "Metric 2 — RMSE in BPM (penalizes outliers)",
    "Per-subject distribution reported (% within ±1, ±3, ±10)",
    "Secondary — SCAMPS synthetic-data baseline",
    "No training step in classical pipeline → whole UBFC = test set",
], size=12)

# Results card
filled_card(s, Inches(6.8), Inches(1.55), Inches(6.0), Inches(2.5),
            fill=GREEN_BG, border=GOOD)
add_text(s, Inches(7.05), Inches(1.7), Inches(5.5), Inches(0.4),
         "EARLY RESULTS (already measured)", size=11, bold=True, color=GOOD)
add_text(s, Inches(7.05), Inches(2.1), Inches(5.5), Inches(0.6),
         "4.06 BPM", size=40, bold=True, color=NAVY)
add_text(s, Inches(7.05), Inches(3.0), Inches(5.5), Inches(0.4),
         "MAE on 42 UBFC real-human subjects",
         size=13, color=SLATE)
add_text(s, Inches(7.05), Inches(3.45), Inches(5.5), Inches(0.4),
         "Rubric target < 10 BPM   ·   2.5× headroom",
         size=11, color=MUTED)

# Comparison table at bottom
add_text(s, Inches(0.5), Inches(4.3), Inches(12.3), Inches(0.4),
         "ACCURACY COMPARISON — measured numbers (docs/evaluation_results.md)",
         size=12, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
rows = [
    ["Dataset",                 "Algorithm",        "n",  "MAE (BPM)", "RMSE (BPM)", "Verdict"],
    ["UBFC-rPPG (real)",       "POS (ours)",       "42", "4.06",      "7.78",       "PASS — rubric < 10"],
    ["UBFC-rPPG (real)",       "CHROM (baseline)", "42", "3.86",      "7.35",       "PASS — tied"],
    ["SCAMPS (synthetic)",     "POS (ours)",       "10", "28.85",     "38.50",      "Synthetic-data limit"],
    ["SCAMPS (synthetic)",     "CHROM (baseline)", "10", "34.83",     "44.05",      "Worse than POS"],
]
add_table(s, Inches(0.5), Inches(4.75),
          [Inches(2.5), Inches(2.6), Inches(0.9), Inches(1.6), Inches(1.7), Inches(3.0)],
          rows, row_h=Inches(0.4), header_h=Inches(0.4), font_size=10)

set_notes(s, """Validation strategy and current measured results. Because the classical pipeline has no learned parameters, there is no training step and the entire 42-subject UBFC dataset functions as the held-out test set. Primary metric is MAE in BPM; secondary is RMSE for outlier sensitivity. Numbers in the bottom table are real measurements from docs/evaluation_results.md.

POS reaches 4.06 BPM MAE on UBFC, clearing the rubric's <10 BPM target by a factor of 2.5. CHROM at 3.86 is statistically tied. On synthetic SCAMPS data the gap widens — POS at 28.85 BPM clearly outperforms CHROM at 34.83 — but synthetic data is a stress test, not the primary benchmark.""")

# ============================================================
# SLIDE 10 — FUTURE MODEL EXTENSIONS + REFERENCES
# ============================================================
s = add_slide()
header(s, "Future model extensions & references",
       "Stretch goals: trained-model artifacts and supporting literature")

# Left — future extensions
filled_card(s, Inches(0.5), Inches(1.55), Inches(6.0), Inches(5.0))
add_text(s, Inches(0.75), Inches(1.7), Inches(5.5), Inches(0.4),
         "PLANNED MODEL EXTENSIONS", size=11, bold=True, color=ACCENT)
add_bullets(s, Inches(0.75), Inches(2.1), Inches(5.5), Inches(4.6), [
    "Level 1 — sklearn regressor over POS-derived features\n"
    "      Real train/test split (30/6/6 by subject); ~1 day effort",
    "Level 2 — 1D CNN on the POS pulse waveform\n"
    "      ~3 days; demonstrates end-to-end DL artifact",
    "Level 3 — PhysNet / TS-CAN from raw frames\n"
    "      ~1 week; GPU required; full deep rPPG model",
    "Level 4 — Blood-pressure regressor on MCD-rPPG\n"
    "      Replaces BP placeholder; requires cuff-paired training data",
    "Bias evaluation — per-skin-tone MAE breakdown on UBFC\n"
    "      Honest reporting of demographic performance gaps",
], size=12)

# Right — references
filled_card(s, Inches(6.8), Inches(1.55), Inches(6.0), Inches(5.0),
            fill=CARD_BG, border=ACCENT)
add_text(s, Inches(7.05), Inches(1.7), Inches(5.5), Inches(0.4),
         "REFERENCES", size=11, bold=True, color=ACCENT)
add_bullets(s, Inches(7.05), Inches(2.1), Inches(5.5), Inches(4.6), [
    "Bawack, R. E., et al. (2021). A framework for understanding "
    "AI research: insights from practice. J. Enterprise Info. Mgmt., 3(2), 645–678.",
    "Wang, W., den Brinker, A. C., Stuijk, S., de Haan, G. (2017). "
    "Algorithmic principles of remote PPG. IEEE TBE 64(7).",
    "de Haan, G., Jeanne, V. (2013). Robust pulse rate from "
    "chrominance-based rPPG. IEEE TBE 60(10).",
    "Liu, X., et al. (2023). rPPG-Toolbox: deep remote PPG toolbox. "
    "NeurIPS Datasets and Benchmarks.",
    "Bobbia, S., et al. (2019). UBFC-rPPG dataset. Pattern Recog. Letters.",
    "McDuff, D., et al. (2022). SCAMPS synthetic dataset. Microsoft Research.",
    "Lugaresi, C., et al. (2019). MediaPipe FaceMesh. arXiv:1906.08172.",
    "Shaffer, F., Ginsberg, J. P. (2017). HRV metrics and norms. "
    "Frontiers in Public Health, 5, 258.",
    "Group 1 Week 3 deliverable — full method-selection rationale.",
], size=10)

set_notes(s, """Closing slide — planned model extensions and references. Four levels of extension scoped on the left, from a 1-day sklearn regressor (which directly answers the train/test/eval question) through a 1-week PhysNet / TS-CAN deep-learning model. Bias evaluation is called out separately because the professor flagged demographic performance as a concern on 2026-05-27.

Right column is the reference list. Bawack et al. is the assignment's supplementary reading; the remaining references are the literature underpinning each algorithm choice in the model. Full method-selection rationale lives in our Week 3 deliverable, which functions as the supporting document to this deck.""")

# Save
prs.save(OUT)
print(f"wrote {OUT}")
print(f"  slides: {len(prs.slides)}")
