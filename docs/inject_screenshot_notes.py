"""Add screenshot-capture instructions to slide 6 of the backend-architecture deck.

Run with:  backend/venv312/bin/python docs/inject_screenshot_notes.py

The script opens the existing pptx, finds slide 6, and writes step-by-step
capture instructions for each of the four screenshot placeholders into that
slide's speaker notes. All other slides and on-slide content are untouched.
"""
from pathlib import Path

from pptx import Presentation

ROOT = Path(__file__).resolve().parent.parent
DECK = ROOT / "docs" / "VitalScan_Group1_Backend_System_Architecture.pptx"

NOTES = """SLIDE 6 — LIVE SITE SCREENSHOTS

How to capture each of the 4 placeholders on this slide. The deck is at
docs/VitalScan_Group1_Backend_System_Architecture.pptx — drop screenshots in
by right-clicking each placeholder → "Change Picture" → "From File".

═══════════════════════════════════════════════════════════════════════
1 · WEBCAM CAPTURE (30-second countdown + camera permission)
═══════════════════════════════════════════════════════════════════════
  a. Open Chrome → https://vitalscan.bkre8tive.com on a Mac/PC with webcam.
  b. Click the "Scan live" button.
  c. Allow camera permission when the browser asks.
  d. The recording UI appears: face preview, corner viewfinder brackets,
     red REC dot top-right, sweeping scan line, countdown ticking down.
  e. Capture mid-recording (so the REC dot + countdown are both visible):
       Cmd + Shift + 4   →   drag a rectangle around the camera modal
     (or hit Space after Cmd+Shift+4 to switch to window-capture mode)
  f. Saves to Desktop as "Screenshot YYYY-MM-DD at HH.MM.SS.png".

═══════════════════════════════════════════════════════════════════════
2 · DASHBOARD (body figure + biomarker cards — the hero shot)
═══════════════════════════════════════════════════════════════════════
  a. After a scan completes, OR just use the persisted state on first load.
  b. Make sure all 4 biomarker cards show real values (HR, HRV, stress, BP)
     and the body-figure markers are colored (green / amber / red).
  c. Capture:
       Cmd + Shift + 4   →   drag a rectangle around the WHOLE main panel:
       body figure on the left + 4 biomarker cards on the right
  d. Frame tight — don't include the bookmarks bar or browser chrome.

═══════════════════════════════════════════════════════════════════════
3 · BIOMARKERS (zoomed-in view of HR / HRV / stress / BP cards)
═══════════════════════════════════════════════════════════════════════
  Two ways to do this — pick whichever looks crisper:
  Option A (preferred — sharper text):
    a. In Chrome, press Cmd + Plus a couple of times to zoom to ~150 %.
    b. Cmd + Shift + 4 → drag a rectangle around just the 4 cards.
    c. Reset zoom with Cmd + 0 when done.
  Option B (faster):
    a. Cmd + Shift + 4 → drag a tight crop around just the 4 cards.
  Make sure the status labels are visible ("Normal range", "Healthy
  range", "Moderate", "Elevated") — those carry the meaning.

═══════════════════════════════════════════════════════════════════════
4 · JSON (/scan response in Chrome DevTools)
═══════════════════════════════════════════════════════════════════════
  a. On vitalscan.bkre8tive.com, open DevTools:  Cmd + Option + I  (or F12)
  b. Click the "Network" tab.
  c. In the filter box, type:  scan
  d. Click "Scan live" or "Upload video" to trigger a scan request.
  e. When the /scan POST completes, click the request row.
  f. Click the "Response" tab in the DevTools right panel — you'll see:
       {
         "biomarkers": {
           "heart_rate": 72,
           "hrv_sdnn": 45,
           "stress_index": 0.6,
           "blood_pressure": { "systolic": 138, "diastolic": 88 }
         }
       }
  g. Capture:
       Cmd + Shift + 4   →   drag a rectangle around the DevTools panel,
       framing BOTH the request row at the top and the JSON response below.
  h. If the JSON is collapsed, click the disclosure triangles to expand it
     before capturing — readability over completeness.

═══════════════════════════════════════════════════════════════════════
INSERTING INTO THIS SLIDE
═══════════════════════════════════════════════════════════════════════
  • Open this deck in PowerPoint or Keynote.
  • Slide 6.
  • For each placeholder: right-click → "Change Picture" → "From File"
    → select the screenshot from your Desktop.
  • Alternative: drag-and-drop the PNG file directly onto the placeholder.
  • The placeholder boxes are pre-sized 6.0 × 2.5 inches — match that
    aspect ratio (~12:5) on your screenshot crops so they don't distort.

═══════════════════════════════════════════════════════════════════════
TALKING POINTS WHILE PRESENTING THIS SLIDE (Germaine)
═══════════════════════════════════════════════════════════════════════
  1. The webcam capture screen — point out the in-browser recording, no
     app to install, 30-second clip, the corner brackets + REC dot are
     UI feedback we built to make the recording state obvious.
  2. The dashboard — body figure with marker positions, four biomarker
     cards showing values pulled live from the backend's JSON response.
  3. The biomarkers up close — call out heart rate in BPM, HRV in ms,
     stress on a 0–1 scale, and BP. Flag that BP is currently mocked
     because classical rPPG can't derive it without cuff calibration.
  4. The /scan JSON — for developers in the room: this is the shared
     API contract any other group can consume. Multipart POST in,
     biomarker JSON out, ~5 seconds end-to-end on real video.

  Close this slide by saying: "the site is live right now — open
  vitalscan.bkre8tive.com on your phone."
"""


def main():
    prs = Presentation(DECK)
    if len(prs.slides) < 6:
        raise SystemExit(f"Expected ≥ 6 slides, found {len(prs.slides)}")
    slide6 = prs.slides[5]  # 0-indexed
    slide6.notes_slide.notes_text_frame.text = NOTES
    prs.save(DECK)
    print(f"wrote screenshot-capture notes into slide 6 of {DECK.name}")
    print(f"  notes length: {len(NOTES)} chars")


if __name__ == "__main__":
    main()
